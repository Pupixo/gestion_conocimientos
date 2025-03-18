from django.shortcuts import render
from rest_framework.decorators import action
from rest_framework.decorators import api_view
from rest_framework_simplejwt.views import TokenObtainPairView
from rest_framework_simplejwt.serializers import TokenObtainPairSerializer
from rest_framework.permissions import IsAuthenticated
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from rest_framework import viewsets

from django.contrib.auth.models import User
from django.db.models import F
from django.db.models import Q
from django.http import HttpResponse, Http404
from django.conf import settings
### Tools
from datetime import datetime, timedelta
# from pytz import timezone
from django.utils import timezone
from django.contrib.auth.hashers import make_password
from django.shortcuts import get_object_or_404
from django.core.files.storage import default_storage
from django.core.files.base import ContentFile
from django.http import JsonResponse
#extras
from django.core.files.base import ContentFile
from django.utils.crypto import get_random_string
from django.http import FileResponse
from django_filters.rest_framework import DjangoFilterBackend

# Modelos Aqui
from apis.models import PerfilUsuario,EstadoUsuario,Roles,Secciones
from apis.models import SeccionesRoles
from apis.models import Cursos,CursoUsuario,Contenidos,Temas,ArchivoTemas,Quiz,PreguntaQuiz,ExamenCurso,PreguntaExamenCurso,NotaExamenCurso,RespuestaUsuarioExamen
from apis.models import Pais,Departamento,Provincia,Distrito,NocsGrupos
from apis.models import LineaTiempo,LineaTiempoEvento

# serializer Aqui
from apis.serializers import UserSerializer,PerfilUsuarioSerializer,RolesSerializer,EstadoUsuarioSerializer,SeccionesSerializer
from apis.serializers import SeccionesRolesSerializer
from apis.serializers import CursosSerializer,ContenidoSerializer,TemasSerializer,ArchivoTemasSerializer,QuizSerializer,PreguntaQuizSerializer,CursoUsuarioMainSerializer,ExamenCursoSerializer,PreguntasExamenCursoSerializer,RespuestaUsuarioExamenSerializer,NotaExamenCursoSerializer,CursoEstadisticaSerializer
from apis.serializers import PaisSerializer,DepartamentoSerializer,ProvinciaSerializer,DistritoSerializer
from apis.serializers import NocsGruposSerializer
from apis.serializers import LineaTiempoSerializer,LineaTiempoEventoSerializer

# filtros
from .filters import PaisFilterSet,DepartamentoFilterSet,ProvinciaFilterSet,DistritoFilterSet
from .filters import NocsGruposFilterSet
from copy import deepcopy

from openpyxl.chart import Reference, Series, BarChart
from openpyxl import Workbook
from openpyxl.styles import PatternFill
from openpyxl.utils.dataframe import dataframe_to_rows

from docx2pdf import convert
from openpyxl.utils import get_column_letter
from adactivexlsx2html import xlsx2html
from reportlab.pdfgen import canvas
from subprocess import Popen
from django.http import QueryDict
from glob import glob

from apis.utils.performance_utils import contar_desempenio
from apis.chatbotfiles.chatbot import get_bot_response

import os
import docx
import openpyxl
import json
import random
import re
import string
import shutil
import pptx 
import pdfplumber
import pdfkit
import mammoth
import io
import tempfile
import subprocess
# from xlsx2html import xlsx2html
# import xlsx2html
import pandas as pd
import xlrd
import xlwt
import pyexcel
import pyexcel_xlsx
import zipfile


class CustomTokenObtainPairSerializer(TokenObtainPairSerializer):
    @classmethod
    def get_token(cls, user):
        token = super().get_token(user)
        if not user.is_active:
            raise serializers.ValidationError("Usuario no activo")

        token["username"] = user.username
        token["modulos"] = ["usuarios", "operaciones", "nodos"]
        return token

class CustomTokenObtainPairView(TokenObtainPairView):
    serializer_class = CustomTokenObtainPairSerializer

class Usuarios(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request,pk=None ,format=None):
        if pk is not None:
            try:
                if pk > 1:
                    usuario = User.objects.get(pk=pk)
                    serializer = UserSerializer(usuario)
                    return Response(serializer.data)
                else:
                    return Response(status=status.HTTP_400_BAD_REQUEST)
            except User.DoesNotExist:
                return Response(status=status.HTTP_404_NOT_FOUND)
        else:
            usuarios = User.objects.all()
            serializer = UserSerializer(usuarios, many=True)
            return Response(serializer.data)

    def post(self, request, format=None):
        password = request.data.get("password")
        mutable_data = deepcopy(request.data)
        pass_hash = make_password(password)
        mutable_data["password"] = pass_hash

        usuario_actual=mutable_data.pop('current_user')
        estado_usu=mutable_data.pop('estado')
        acronimo_rol=mutable_data.pop('acronimo')

        usuarios_instance = User.objects.filter(username=mutable_data["username"]).first()

        if usuarios_instance:
            roles_instance = Roles.objects.filter(acronimo=acronimo_rol[0]).first()
            if roles_instance:
                esta_usu_instance = EstadoUsuario.objects.filter(estado_usu_name=estado_usu[0]).first()
                if esta_usu_instance:
                        perfil_user_instance = PerfilUsuario.objects.filter(user= usuarios_instance.id ).first()
                        perfil_user_instance.rol =roles_instance
                        perfil_user_instance.estado =esta_usu_instance
                        perfil_user_instance.save()
                        return Response({"msg": "¡Se actualizo al usuario!"}, status=status.HTTP_200_OK)
                else:
                    return Response({"msg": "estado "+ estado_usu[0]+ " no existe" }, status=status.HTTP_400_BAD_REQUEST)
            else:
                return Response({"msg": "acronimo rol "+ acronimo_rol[0]+ " no existe"}, status=status.HTTP_400_BAD_REQUEST)

        else:
            serializer = UserSerializer(data=mutable_data)
            if serializer.is_valid():
                serializer.save()

                pk = serializer.data["id"]
                roles_instance = Roles.objects.filter(acronimo=acronimo_rol[0]).first()

                if roles_instance:

                    esta_usu_instance = EstadoUsuario.objects.filter(estado_usu_name=estado_usu[0]).first()
                    if esta_usu_instance:
                        datos={ "user":pk,"rol":roles_instance.id,"estado":esta_usu_instance.id,'usuario_reg':usuario_actual[0] }
                        serializer_data = PerfilUsuarioSerializer(data=datos)
                        if serializer_data.is_valid():
                            serializer_data.save()
                            return Response({"msg": "¡Se registro al usuario!"}, status=status.HTTP_200_OK)
                        else:
                            return Response({"msg": "ERROR", "detail": serializer_data.errors}, status=status.HTTP_400_BAD_REQUEST)
                    else:
                        return Response({"msg": "estado "+ estado_usu[0]+ " no existe", "detail":"No existe "+acronimo_rol[0] }, status=status.HTTP_400_BAD_REQUEST)
                    
                else:
                    return Response({"msg": "acronimo rol "+ acronimo_rol[0]+ " no existe", "detail":"No existe "+acronimo_rol[0] }, status=status.HTTP_400_BAD_REQUEST)
           
            else:
                return Response({"msg": "ERROR", "detail": serializer.errors}, status=status.HTTP_400_BAD_REQUEST)


    def put(self, request, pk, format=None):
        try:
            if pk > 1:
                user = User.objects.get(pk=pk)
            else:
                return Response(status=status.HTTP_400_BAD_REQUEST)
        except User.DoesNotExist:
            return Response(status=status.HTTP_400_BAD_REQUEST)

        data = request.data.dict()
        data["password"] = user.password
        serializer = UserSerializer(user, data=data)
        if serializer.is_valid():
            serializer.save()
            return Response(serializer.data, status=status.HTTP_201_CREATED)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

class UsuariosIndividual(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, format=None):
        password = request.data.get("password")
        mutable_data = deepcopy(request.data)
        pass_hash = make_password(password)
        mutable_data["password"] = pass_hash
        fecha_creacion = datetime.now()
        is_superuser =False

        if mutable_data["rol"] ==3 :
            is_superuser =True

        usuarios_instance = User.objects.filter(username=mutable_data["username"]).first()
        if usuarios_instance:
            return Response({"msg": "¡El usuario ya existe!"}, status=status.HTTP_200_OK)
        else:

            datos={'username': mutable_data["username"], 'password': mutable_data["password"], 
                    'first_name': mutable_data["first_name"] , 'last_name':mutable_data["last_name"],
                    'email': mutable_data["email"], 'is_staff':mutable_data["is_staff"], 'is_active': mutable_data["is_active"],
                    'date_joined': fecha_creacion,'is_superuser':is_superuser}

            serializer = UserSerializer(data=mutable_data)
            if serializer.is_valid():
                serializer.save()
                pk = serializer.data["id"]

                datos={ "user":pk,"rol":mutable_data["rol"], "estado":mutable_data["estado"],'estado_fecha':fecha_creacion,'usuario_act_estado':mutable_data["current_user"]}
                serializer_data = PerfilUsuarioSerializer(data=datos)
                if serializer_data.is_valid():
                    serializer_data.save()
                    return Response({"msg": "¡Se registro al usuario!"}, status=status.HTTP_200_OK)
                else:
                    return Response({"msg": "Error al registrar perfil usuario", "detail": serializer_data.errors}, status=status.HTTP_400_BAD_REQUEST)
            else:
                return Response({"msg": "Error al registrar usuario", "detail": serializer.errors}, status=status.HTTP_400_BAD_REQUEST)


    def put(self, request, pk, format=None):
        try:
            if pk >= 1:
                user = User.objects.get(pk=pk)
            else:
                return Response(status=status.HTTP_400_BAD_REQUEST)
        except User.DoesNotExist:
            return Response(status=status.HTTP_400_BAD_REQUEST)

        data = request.data.dict()
        data["password"] = make_password(data["password"])

        if data["password"] == '':
            del data["password"]

        estado=data.pop('estado')
        rol=data.pop('rol')
        usuario_actual=data.pop('current_user')

        serializer = UserSerializer(user, data=data)
        if serializer.is_valid():
            serializer.save()
            fecha_creacion = datetime.now()

            perfil_user_instance = PerfilUsuario.objects.filter(user= pk ).first()
            esta_usu_instance = EstadoUsuario.objects.filter(id=estado).first()

            if perfil_user_instance.estado is not None:
                estado_int = int(perfil_user_instance.estado.id)
                if estado_int != int(estado):
                    perfil_user_instance.estado =esta_usu_instance
                    perfil_user_instance.estado_fecha =fecha_creacion
                    perfil_user_instance.usuario_act_estado =usuario_actual
            else:
                perfil_user_instance.estado =esta_usu_instance
                perfil_user_instance.estado_fecha =fecha_creacion
                perfil_user_instance.usuario_act_estado =usuario_actual


            roles_instance = Roles.objects.filter(id=rol).first()
            perfil_user_instance.rol =roles_instance
            perfil_user_instance.usuario_act =usuario_actual
            perfil_user_instance.editar_fecha =fecha_creacion

            perfil_user_instance.save()

            return Response(serializer.data, status=status.HTTP_201_CREATED)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)


    def delete(self, request, pk, format=None):
        try:
            if pk > 1:
                user = User.objects.get(pk=pk)
            else:
                return Response(status=status.HTTP_400_BAD_REQUEST)
        except User.DoesNotExist:
            return Response(status=status.HTTP_400_BAD_REQUEST)

        user.is_active =False
        user.save()
        # user.delete()
        return Response(status=status.HTTP_204_NO_CONTENT)

class UsuariosIndividualFoto(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, format=None):

        upload_dir = os.path.join(settings.MEDIA_ROOT, 'usuarios/perfil_fotos')
        data = request.data
        foto = request.FILES.get("foto")
        id_usuario = data.get('current_user')
        print("foto ..............................",foto)
        if foto:
            # Extract file extension
            extension = foto.content_type.split('/')[1]
            # Generate unique filename
            filename = f'yn_{get_random_string(length=10)}.{extension}'
            while os.path.exists(os.path.join(upload_dir, filename)):
                filename = f'yn_{get_random_string(length=10)}.{extension}'
            # Create directory with user ID
            user_dir = os.path.join(upload_dir, str(id_usuario))
            os.makedirs(user_dir, exist_ok=True)  # Create if doesn't exist

            print("user_dir............................",user_dir)
            # Save the image inside user directory
            full_path = os.path.join(user_dir, filename)
            with open(full_path, 'wb') as destination:
                for chunk in foto.chunks():
                    destination.write(chunk)
            # ... (consider saving metadata to database if needed)
            perfil_user_instance = PerfilUsuario.objects.filter(user= id_usuario ).first()
            #  ... Eliminar archivo anterior si es que habia 
            if perfil_user_instance.foto is  None or perfil_user_instance.foto == '':
                print("nada se hace")
            else :
                ruta_archivo_anterior =upload_dir+'/'+id_usuario+'/'+perfil_user_instance.foto
                if os.path.exists(ruta_archivo_anterior):
                    os.remove(ruta_archivo_anterior)
                    print("elimino.......")
                else:
                    print("no elimino.......")
                    # La imagen no existe, no se hace nada
                    # pass

            perfil_user_instance.foto = id_usuario+ '/' +filename
            perfil_user_instance.save()
            # Return success response (replace with your desired response)
            return Response({'message': 'Image uploaded successfully!'}, status=status.HTTP_201_CREATED)
        else:
            # Handle case where no file is uploaded
            return Response({'error': 'No image file uploaded!'}, status=400)


class UsuariosIndividualDatos(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, format=None):
        data = request.data
        id_usuario = data.get('current_user')
        print("data ..............................",data)
        if data:
            # Extract file extension
           
            return Response({'message': 'Image uploaded successfully!'}, status=status.HTTP_201_CREATED)
        else:
            # Handle case where no file is uploaded
            return Response({'error': 'No image file uploaded!'}, status=400)



# localidades vistas
class PaisesViewSet(viewsets.ModelViewSet):
    queryset = Pais.objects.all()
    serializer_class = PaisSerializer
    filter_backends = [DjangoFilterBackend]
    filterset_class = PaisFilterSet
    permission_classes = [IsAuthenticated]

    def list(self, request, *args, **kwargs):
        # Usa el método base `list` para obtener el queryset filtrado
        response = super().list(request, *args, **kwargs)
        # Aquí puedes añadir lógica adicional si es necesario
        return response

class DepartamentoViewSet(viewsets.ModelViewSet):
    queryset = Departamento.objects.all()
    serializer_class = DepartamentoSerializer
    filter_backends = [DjangoFilterBackend]
    filterset_class = DepartamentoFilterSet
    permission_classes = [IsAuthenticated]

    def list(self, request, *args, **kwargs):
        # Usa el método base `list` para obtener el queryset filtrado
        response = super().list(request, *args, **kwargs)
        # Aquí puedes añadir lógica adicional si es necesario
        return response
    

    @action(detail=False, methods=['get'])
    def list_by_type(self, request, id, type=None):
        if type == 'usuario':
            archivos = self.queryset.filter(usuarios=id,status=True)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)
    
    
class ProvinciaViewSet(viewsets.ModelViewSet):
    queryset = Provincia.objects.all()
    serializer_class = ProvinciaSerializer
    filter_backends = [DjangoFilterBackend]
    filterset_class = ProvinciaFilterSet
    permission_classes = [IsAuthenticated]

    def list(self, request, *args, **kwargs):
        # Usa el método base `list` para obtener el queryset filtrado
        response = super().list(request, *args, **kwargs)
        # Aquí puedes añadir lógica adicional si es necesario
        return response

    @action(detail=False, methods=['get'])
    def list_by_type(self, request, id, type=None):
        if type == 'usuario':
            archivos = self.queryset.filter(usuarios=id,status=True)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)
    
    
class DistritoViewSet(viewsets.ModelViewSet):
    queryset = Distrito.objects.all()
    serializer_class = DistritoSerializer
    filter_backends = [DjangoFilterBackend]
    filterset_class = DistritoFilterSet
    permission_classes = [IsAuthenticated]

    def list(self, request, *args, **kwargs):
        # Usa el método base `list` para obtener el queryset filtrado
        response = super().list(request, *args, **kwargs)
        # Aquí puedes añadir lógica adicional si es necesario
        return response

    @action(detail=False, methods=['get'])
    def list_by_type(self, request, id, type=None):
        if type == 'usuario':
            archivos = self.queryset.filter(usuarios=id,status=True)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)
    
        
class NocsGruposViewSet(viewsets.ModelViewSet):
    queryset = NocsGrupos.objects.all()
    serializer_class = NocsGruposSerializer
    filter_backends = [DjangoFilterBackend]
    filterset_class = NocsGruposFilterSet
    permission_classes = [IsAuthenticated]

    def list(self, request, *args, **kwargs):
        # Usa el método base `list` para obtener el queryset filtrado
        response = super().list(request, *args, **kwargs)
        # Aquí puedes añadir lógica adicional si es necesario
        return response


class UsuariosInfo(APIView):

    permission_classes = [IsAuthenticated]
    def post(self, request, format=None):

        datos=request.data
        all_values = list(datos.values())
        user_crt = all_values[0]
        pass_crt = all_values[1]

        data_output = (
                        'id',
                        'password',
                        'last_login',
                        'is_superuser',
                        'username',
                        'first_name',
                        'last_name',
                        'email',
                        'is_staff',
                        'is_active',
                        'date_joined',
        )

        user_data = User.objects.filter(username = user_crt).first()
        user_data = PerfilUsuario.objects.filter(user = user_data.id)
        serializer = PerfilUsuarioSerializer(user_data, many=True)
        return Response(serializer.data)  

class UsuariosPassword(APIView):
    
    permission_classes = [IsAuthenticated]

    def put(self, request, pk, format=None):

        try:
            user = User.objects.get(pk=pk)
        except User.DoesNotExist:
            return Response(status=status.HTTP_400_BAD_REQUEST)

        new_password = request.data.dict()
        user.set_password(new_password["password"])
        user.save()
        return Response(status=status.HTTP_204_NO_CONTENT)

#rol
class RolAnalisis(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request,pk=None ,format=None):
        if pk is not None:
            try:
                if pk > 1:
                    usuario = Roles.objects.get(pk=pk)
                    serializer = RolesSerializer(usuario)
                    return Response(serializer.data)
                else:
                    return Response(status=status.HTTP_400_BAD_REQUEST)
            except Roles.DoesNotExist:
                return Response(status=status.HTTP_404_NOT_FOUND)
        else:
            usuarios = Roles.objects.filter(status=True)
            serializer = RolesSerializer(usuarios, many=True)
            return Response(serializer.data)
        

    def post(self, request, format=None):

        # return Response({"msg": "¡Se registro el rol!"}, status=status.HTTP_200_OK)

        # datos={ "user":pk,"rol":mutable_data["rol"], "estado":mutable_data["estado"],'estado_fecha':fecha_creacion,'usuario_act_estado':mutable_data["current_user"]}
        datos=request.data
        serializer = RolesSerializer(data=datos)
        if serializer.is_valid():
            serializer.save()
            pk = serializer.data["id"]
            return Response({"msg": "¡Se registro el rol!"}, status=status.HTTP_200_OK)
        else:
            return Response({"msg": "Error al registrar al rol", "detail": serializer.errors}, status=status.HTTP_400_BAD_REQUEST)

    def put(self, request, pk, format=None):
        try:
            if pk > 1:
                rol = Roles.objects.get(pk=pk)
            else:
                return Response(status=status.HTTP_400_BAD_REQUEST)
        except Roles.DoesNotExist:
            return Response(status=status.HTTP_400_BAD_REQUEST)

        serializer = RolesSerializer(rol, data=request.data)
        if serializer.is_valid():
            serializer.save()
            return Response(serializer.data, status=status.HTTP_201_CREATED)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

#secciones
class SeccionAnalisis(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request,pk=None ,format=None):

        padre = request.query_params["id_padre"]
        tipo = request.query_params["tipo"]

        if pk is not None:
            try:
                if pk > 1:
                    usuario = Secciones.objects.get(pk=pk)
                    serializer = SeccionesSerializer(usuario)
                    return Response(serializer.data)
                else:
                    return Response(status=status.HTTP_400_BAD_REQUEST)
            except Secciones.DoesNotExist:
                return Response(status=status.HTTP_404_NOT_FOUND)
        else:
            # secciones
            if tipo == '1':
                secciones = Secciones.objects.filter(status=True,id_padre=None)
            serializer = SeccionesSerializer(secciones, many=True)
            return Response(serializer.data)
        
#secciones roles
class SeccionRolesAnalisis(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request,pk=None ,format=None):

        seccion = request.query_params["seccion"]
        tipo = request.query_params["tipo"]
        rol_id = request.query_params["rol_id"]

        # secciones
        if tipo == '1':
            secciones = SeccionesRoles.objects.filter(role=rol_id,status_selec=True)
        
        serializer = SeccionesRolesSerializer(secciones, many=True)
        return Response(serializer.data)
    
    def post(self, request, format=None):

        lista =request.data.get("listaSeleccionados")
        detalle = json.loads(lista)
        SeccionesRoles.objects.filter(role=request.data.get("role")).update(status_selec=False)
        for pk in detalle:
            seccion = SeccionesRoles.objects.filter(role=request.data.get("role"),seccion=pk["id"])
            if  seccion.count() == 0:
                datos={ 
                    "seccion": pk["id"],
                    "role":request.data.get("role"),
                    "detalle":pk["detalle"],
                    "status_selec":True
                }
                serializer = SeccionesRolesSerializer(data=datos)
                if serializer.is_valid():
                    serializer.save()
                else:
                    return Response({"msg": "Error al registrar la secciòn", "detail": serializer.errors}, status=status.HTTP_400_BAD_REQUEST)             
            else :
                SeccionesRoles.objects.filter(role=request.data.get("role"),seccion=pk["id"]).update(status_selec=True)
                # seccion = SeccionesRoles.objects.get(id=pk["id_secc_rol"])
                # seccion.status_selec=True
      
        return Response({"msg": "¡Se hecho!"}, status=status.HTTP_200_OK)

#est
class estadoUsuarioAnalisis(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request,pk=None ,format=None):
        if pk is not None:
            try:
                if pk > 1:
                    usuario = EstadoUsuario.objects.get(pk=pk)
                    serializer = EstadoUsuarioSerializer(usuario)
                    return Response(serializer.data)
                else:
                    return Response(status=status.HTTP_400_BAD_REQUEST)
            except EstadoUsuario.DoesNotExist:
                return Response(status=status.HTTP_404_NOT_FOUND)
        else:
            usuarios = EstadoUsuario.objects.filter(status=True)
            serializer = EstadoUsuarioSerializer(usuarios, many=True)
            return Response(serializer.data)
       
class BusquedaArchivo(APIView):
    permission_classes = [IsAuthenticated]
       
    def search_pdf(self,file_path, query):
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_without_spaces = query.replace(" ", "")
                    if text_without_spaces.lower() in text.lower():
                        return True
        return False

    def search_word(self, file_path, query):
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            if query.lower() in para.text.lower():
                return True
    
        # Buscar en las tablas del documento
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if query.lower() in cell.text.lower():
                        return True
    
        return False

    def search_excel(self, file_path, query):
        wb = openpyxl.load_workbook(file_path, read_only=True)
        valor_buscado_str = str(query).lower()

        for sheet in wb.sheetnames:
            ws = wb[sheet]          
            # Iterar a través de las filas
            for row in ws.rows:
                # Iterar a través de las celdas en la fila
                for cell in row:
                    # Convertir el valor de la celda a cadena de texto
                    cell_str = str(cell.value).lower()
                    # Comprobar si el texto o el número buscado coincide
                    if valor_buscado_str in cell_str:
                        return True
            

        #  # Buscar en los gráficos
        # for chart in ws._charts:
        #     for series in chart.series:
        #         for point in series.points:
        #             if query.lower() in str(point).lower():
        #                 return True
                        
        return False
    
    def search_ppt(self, file_path, query):
        # prs = pptx.Presentation(file_path)
        # for slide in prs.slides:
        #     for shape in slide.shapes:
        #         if hasattr(shape, "text") and query.lower() in shape.text.lower():
        #             return True
        return False
 
    def search_txt(self, file_path, query):   
        matching_lines = []
        try:
            with open(file_path, 'r') as file:
                for line in file:
                    if query in line:
                        matching_lines.append(line.strip())
                        return True

        except FileNotFoundError:
            print(f"Error: File not found: {file_path}")
        return False

    def search_in_file(self, file_path, query):
        try:
            if file_path.endswith('.docx') or file_path.endswith('.doc')  :
                return self.search_word(file_path, query)
            elif file_path.endswith('.xlsx') or file_path.endswith('.xls'):
                return self.search_excel(file_path, query)
            elif file_path.endswith('.pptx') or file_path.endswith('.ppt') :
                return self.search_ppt(file_path, query)
            elif file_path.endswith('.pdf'):
                return self.search_pdf(file_path, query)
            elif file_path.endswith('.txt'):
                return self.search_txt(file_path, query)
        except Exception as e:
            print(f"Error processing file {file_path}: {e}")
        return False

    def get(self, request, query=None, format=None):
        if query is None:
            return Response([])  # O podrías devolver un mensaje de error.

        # media_dir = os.path.join(settings.MEDIA_ROOT, 'archivos')
        # files_matched = []
        # for file_name in os.listdir(media_dir):
        #     full_path = os.path.join(media_dir, file_name)

        #     if query.lower() in file_name.lower():
        #         files_matched.append(file_name)
        #     elif self.search_in_file(full_path, query):
        #         files_matched.append(file_name)
        # return Response(files_matched)



        media_dir = os.path.join(settings.MEDIA_ROOT, 'archivos')
        files_matched = []

        # Recorrer todas las carpetas y subcarpetas de 'media_dir'
        for root, _, files in os.walk(media_dir):
            for file_name in files:
                full_path = os.path.join(root, file_name)  # Ruta completa del archivo

                # Comprobar si el archivo contiene la cadena 'query' (insensible a mayúsculas/minúsculas)
                if query.lower() in file_name.lower():
                    files_matched.append(full_path)  # Agregar la ruta completa
                elif self.search_in_file(full_path, query):
                    files_matched.append(full_path)  # Agregar si el contenido del archivo contiene la cadena

        return Response(files_matched) 
    
    #visualizar archivos

    def convert_word(self, file_path):
        folderDir = os.path.join(settings.MEDIA_ROOT, 'archivospdf')

        # convert(file_path,folderDir+ '/sample.pdf')
        with open(file_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html = result.value 
            # messages = result.messages
        # f = open(folderDir + '/sample.html',"w")
        # f.write(html)
        # f.close()
        return html
        
    def convert_excel(self, file_path):
        media_dir = os.path.join(settings.MEDIA_ROOT, 'archivospdf')
        # xlsx2html(file_path, media_dir+'output.html')
        # self.convert_xlsx_to_html_files(file_path,media_dir)
        # xlsx_file = open(file_path, 'rb')
        # out_file = io.StringIO()
        # xlsx2html(xlsx_file, out_file, locale='en')
        # out_file.seek(0)
        # # out_file.getvalue()
        # result_html = out_file.read()
        # return result_html
    

        xlsx_file = open(file_path, 'rb')
        out_file = io.StringIO()
        xlsx2html(xlsx_file, out_file, locale='en')
        out_file.seek(0)
        result_html = out_file.read()

        # out_stream = xlsx2htmladapt('path/to/example.xlsx')
        # out_stream.seek(0)
        # print(out_stream.read())
        return result_html

    def convert_ppt(self, file_path):
        contents = io.StringIO()
        # contents = open"C:\\Users\\Suleiman JK\\Desktop\\Static_hash\\test","r")
        with open("suleiman.html", "w") as e:
            for lines in contents.readlines():
                e.write("<pre>" + lines + "</pre> <br>\n")
        return contents

    def convert_txt_to_html(self,file_path):
        try:
            # Read the contents of the text file into a StringIO object
            with open(file_path, "r") as txt_file:
                # Create a string buffer to store the HTML content
                html_content = io.StringIO()
                # Write a simple HTML header
                html_content.write("<html><head><title>Contenido del txt</title></head><body>")
                # Read each line from the text file, wrapping in <pre> to preserve formatting
                for line in txt_file:
                    html_content.write(f"<pre>{line.strip()}</pre><br>")
                # Write a simple HTML footer
                html_content.write("</body></html>")
            # Return the entire HTML content as a string
            return html_content.getvalue()
        
        except Exception as e:
            print(f"Error converting text to HTML: {e}")
            return None

    def archivo_generico_preview(self, file_path):
        return 'todo correcto'

    def convert_file(self, file_path):
        try:
            if file_path.endswith('.docx') or file_path.endswith('.doc'):
                return self.convert_word(file_path)
            elif file_path.endswith('.xlsx') or file_path.endswith('.xls'):
                return self.convert_excel(file_path)
            elif file_path.endswith('.pptx') or file_path.endswith('.ppt'):
                return self.convert_ppt(file_path)
            elif file_path.endswith('.pdf'):
                return self.search_pdf(file_path)
            elif file_path.endswith('.txt'):
                return self.convert_txt_to_html(file_path)
            else:
                return self.archivo_generico_preview(file_path)

        except Exception as e:
            print(f"Error processing file {file_path}: {e}")
        return False


    def post(self, request, format=None):
        datos=request.data
        print("datos..........................",datos)
        print("datos..........................",datos['docfile'])        
        media_dir = os.path.join(settings.MEDIA_ROOT, 'archivos')

        file_name = datos['docfile']
        doc = media_dir +'/'+ file_name
        print("doc.................",doc)

        html = self.convert_file(doc)

        return Response({"html": html}, status=status.HTTP_200_OK)



class RutasArchivo(APIView):
    permission_classes = [IsAuthenticated]
    # Helper function to get folder structure recursively
    def get_folder_structure(self, path, base_path):
        estructura_directorio = []

        carpeta_relativa = os.path.relpath(path, base_path)
        # Listar y ordenar las subcarpetas y archivos
        sub_carpetas = sorted([d for d in os.listdir(path) if os.path.isdir(os.path.join(path, d))])  # Para subcarpetas
        archivos = sorted([f for f in os.listdir(path) if os.path.isfile(os.path.join(path, f))])  # Para archivos

        # Obtener la ruta relativa y la absoluta
        carpeta_relativa = os.path.relpath(path, base_path)

        carpeta_absoluta = os.path.abspath(path)
        # El prefijo a eliminar
        prefijo = "/usr/src/app/media/"
        # Eliminar el prefijo si está presente en el string
        ruta_sin_prefijo = carpeta_absoluta.replace(prefijo, "")
        string_modificado = ruta_sin_prefijo.replace("/", "$")

        contenido = {
            'carpeta': carpeta_relativa,
            'sub_carpetas': sub_carpetas,
            'archivos': archivos,
            'ruta': string_modificado
        }

        # Agregar el contenido al array principal
        estructura_directorio.append(contenido)

        return estructura_directorio

    def get(self, request, pk=None, format=None):
        # Directorio raíz
        media_dir = os.path.join(settings.MEDIA_ROOT, 'archivos')
        if pk is not None:

            media_dir = os.path.join(settings.MEDIA_ROOT)
            string_modificado = pk.replace("$", "/")
            
            ruta = media_dir+'/' + string_modificado
            print("ruta...........................",ruta)
            data_final = self.get_folder_structure(ruta, media_dir)
        else:
            # Obtener la estructura del directorio
            data_final = self.get_folder_structure(media_dir, media_dir)
            # Retornar la respuesta con la estructura del directorio
        
        return Response(data_final)



class CrearCarpetas(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        datos = request.data

        # Validar si 'rutafolder' y 'newfolder' existen en los datos
        rutafolder = datos.get('rutafolder')
        newfolder = datos.get('newfolder')

        if not rutafolder or not newfolder:
            return Response({"error": "Los campos 'rutafolder' y 'newfolder' son requeridos"}, 
                            status=status.HTTP_400_BAD_REQUEST)

        # Sanitizar nombres de carpeta (opcional, según sea necesario)
        rutafolder = os.path.normpath(rutafolder).strip('/')
        newfolder = os.path.normpath(newfolder).strip('/')

        print("rutafolder.......................................",rutafolder)
        print("newfolder......................................",newfolder)

        # Ruta base donde se creará la carpeta
        base_path = settings.MEDIA_ROOT  # Cambia esta ruta según tu entorno
        folder_path = os.path.join(base_path, rutafolder, newfolder)
        print("folder_path......................................",folder_path)

        try:
            # Crear la carpeta si no existe
            if not os.path.exists(folder_path):
                os.makedirs(folder_path)
                return Response({"message": "Carpeta creada correctamente"}, status=status.HTTP_201_CREATED)
            else:
                return Response({"message": "La carpeta ya existe"}, status=status.HTTP_200_OK)
        
        except OSError as e:
            return Response({"error": f"No se pudo crear la carpeta: {str(e)}"}, 
                            status=status.HTTP_500_INTERNAL_SERVER_ERROR)



class CrearArchivo(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        datos = request.data

        # Validar si 'rutafolder', 'newfolder' y 'archivo_file' existen en los datos
        rutafolder = datos.get('rutafolder')
        nomfile = datos.get('nom_file')
        archivo_file = request.FILES.get('archivo')  # Aquí obtenemos el archivo

        if not rutafolder or not nomfile or not archivo_file:
            return Response({"error": "Los campos 'rutafolder', 'nomfile' y 'archivo' son requeridos"}, 
                            status=status.HTTP_400_BAD_REQUEST)

        # Sanitizar nombres de carpeta (opcional, según sea necesario)
        rutafolder = os.path.normpath(rutafolder).strip('/')
        nomfile = os.path.normpath(nomfile).strip('/')

        print("rutafolder.......................................", rutafolder)
        print("nomfile......................................", nomfile)

        # Ruta base donde se creará la carpeta
        base_path = settings.MEDIA_ROOT  # Cambia esta ruta según tu entorno
        folder_path = os.path.join(base_path, rutafolder)
        print("folder_path......................................", folder_path)

        try:
            # Creamos la carpeta si no existe
            os.makedirs(folder_path, exist_ok=True)

            # Guardamos el archivo en la carpeta correspondiente
            file_path = os.path.join(folder_path, nomfile)
            
            with open(file_path, 'wb+') as destination:
                for chunk in archivo_file.chunks():
                    destination.write(chunk)
            return Response({"message": "Archivo subido correctamente"}, status=status.HTTP_201_CREATED)
        
        except OSError as e:
            error_msg = f"Error al crear la carpeta o guardar el archivo: {str(e)}"
            return Response({"error": error_msg}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)




class EliminoArchivo(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        datos = request.data
        rutafolder = datos.get('ruta')

        if not rutafolder:
            return Response({"error": "Los campos 'rutafolder' son requeridos"}, 
                            status=status.HTTP_400_BAD_REQUEST)

        # Sanitizar nombres de carpeta (opcional, según sea necesario)
        rutafolder = os.path.normpath(rutafolder).strip('/')
        print("--------------------------------------------------------------")
        print("---------------------eliminar archivo-------------------------")
        print("--------------------------------------------------------------")
        print("rutafolder.......................................", rutafolder)
        # Ruta base donde se creará la carpeta
        base_path = settings.MEDIA_ROOT  # Cambia esta ruta según tu entorno
        folder_path = os.path.join(base_path, rutafolder)
        print("folder_path......................................", folder_path)
        newfolder = os.path.normpath(folder_path).strip('/')

        # Verificar si el archivo existe antes de eliminarlo
        if os.path.exists(newfolder):
            os.remove(newfolder)
            print(f"Archivo eliminado: {newfolder}")
            return Response({"message": "Archivo eliminado correctamente"}, status=status.HTTP_201_CREATED)

        else:
            error_msg = f"Error al eliminar el archivo"
            print(f"El archivo no existe: {newfolder}")
            return Response({"error": error_msg}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)




###############################################################################################################################################################################
###############################################################################################################################################################################

#cursos

class FuncionesBasicasVista:
     
    def calcular_puntaje(self, opciones, lista_indices_opciones_correctas):

        # Obtener todos los índices que tienen valor True en opciones
        indices_true_usu = [i for i, valor in enumerate(opciones) if valor]
        print("indices_true_usu .....", indices_true_usu)        
        print("lista_indices_opciones_correctas .....", lista_indices_opciones_correctas)

        print("set(lista_indices_opciones_correctas) .....", set(lista_indices_opciones_correctas))
        print("set(indices_true_usu) .....", set(indices_true_usu))

        # Verificar si todos los índices en lista_indices_opciones_correctas son correctos
        if set(lista_indices_opciones_correctas) == set(indices_true_usu):
            return 4
            # return True

        else:
            return 0
            # return False



    def convertir_a_booleano(self,valor):
        if valor.lower() == 'true':
            return True
        elif valor.lower() == 'false':
            return False
        else:
            raise ValueError(f'“{valor}” value must be either "true" or "false".')

class FileHandler:

    def guardar_archivo(self, archivo, id_curso, tipo,curso):
        if not archivo:
            print(f"No se encontró el archivo de tipo {tipo}.")
            return
        try:
            # Extraer la extensión del archivo
            extension = archivo.name.split('.')[-1]
            # Generar nombre de archivo único
            filename = f'{tipo}_{curso}_{id_curso}.{extension}'
            # Definir la ruta base y completa para la carpeta y archivo
            ruta_base = os.path.join(settings.MEDIA_ROOT, "cursos", tipo, f"curso_{id_curso}")
            if not os.path.exists(ruta_base):
                os.makedirs(ruta_base)
                print(f"Carpeta creada: '{ruta_base}'")
            else:
                print(f"La carpeta ya existe: '{ruta_base}'")
            # Ruta completa del archivo
            full_path = os.path.join(ruta_base, filename)
            # Guardar el archivo en el directorio especificado
            with open(full_path, 'wb') as destination:
                for chunk in archivo.chunks():
                    destination.write(chunk)
            print(f"Archivo de tipo {tipo} guardado en: '{full_path}'")
            return full_path
        except Exception as e:
            print(f"Error al guardar el archivo de tipo {tipo}: {e}")
            return None

    def guardar_archivo_temas(self, archivo, id_contenido, tipo):
        if not archivo:
            print(f"No se encontró el archivo de tipo {tipo}.")
            return
        try:
            # Obtener el nombre original del archivo
            # filename = 'fondo_tema_'+id_contenido;
            filename_original = archivo.name
            filename = filename_original.replace(' ', '_')
            # Definir la ruta base y completa para la carpeta y archivo
            ruta_base = os.path.join(settings.MEDIA_ROOT, "cursos", "contenido", "tema",tipo, f"tema_{id_contenido}")
            if not os.path.exists(ruta_base):
                os.makedirs(ruta_base)
                print(f"Carpeta creada: '{ruta_base}'")
            else:
                print(f"La carpeta ya existe: '{ruta_base}'")
            # Ruta completa del archivo
            full_path = os.path.join(ruta_base, filename)
            # Guardar el archivo en el directorio especificado
            with open(full_path, 'wb') as destination:
                for chunk in archivo.chunks():
                    destination.write(chunk)
            print(f"Archivo de tipo {tipo} guardado en: '{full_path}'")
            return full_path
        except Exception as e:
            print(f"Error al guardar el archivo de tipo {tipo}: {e}")
            return None
        
    def guardar_archivo_quiz(self, archivo, id_curso, tipo):
        if not archivo:
            print(f"No se encontró el archivo de tipo {tipo}.")
            return
        try:
            # Extraer la extensión del archivo
            extension = archivo.name.split('.')[-1]
            # Generar nombre de archivo único
            filename = f'{tipo}_quiz_{id_curso}.{extension}'
            # Definir la ruta base y completa para la carpeta y archivo
            ruta_base = os.path.join(settings.MEDIA_ROOT, "cursos", "contenido", "quiz", tipo, f"curso_{id_curso}")
            if not os.path.exists(ruta_base):
                os.makedirs(ruta_base)
                print(f"Carpeta creada: '{ruta_base}'")
            else:
                print(f"La carpeta ya existe: '{ruta_base}'")
            # Ruta completa del archivo
            full_path = os.path.join(ruta_base, filename)
            # Guardar el archivo en el directorio especificado
            with open(full_path, 'wb') as destination:
                for chunk in archivo.chunks():
                    destination.write(chunk)
            print(f"Archivo de tipo {tipo} guardado en: '{full_path}'")
            return full_path
        except Exception as e:
            print(f"Error al guardar el archivo de tipo {tipo}: {e}")
            return None
        
    def guardar_archivo_preguntas(self, archivo, id_curso, tipo):
        if not archivo:
            print(f"No se encontró el archivo de tipo {tipo}.")
            return
        try:
            # Extraer la extensión del archivo
            extension = archivo.name.split('.')[-1]
            # Generar nombre de archivo único
            filename = f'{tipo}_quiz_{id_curso}.{extension}'
            # Definir la ruta base y completa para la carpeta y archivo
            ruta_base = os.path.join(settings.MEDIA_ROOT, "cursos", "contenido", "quiz", tipo, f"curso_{id_curso}")
            if not os.path.exists(ruta_base):
                os.makedirs(ruta_base)
                print(f"Carpeta creada: '{ruta_base}'")
            else:
                print(f"La carpeta ya existe: '{ruta_base}'")
            # Ruta completa del archivo
            full_path = os.path.join(ruta_base, filename)
            # Guardar el archivo en el directorio especificado
            with open(full_path, 'wb') as destination:
                for chunk in archivo.chunks():
                    destination.write(chunk)
            print(f"Archivo de tipo {tipo} guardado en: '{full_path}'")
            return full_path
        except Exception as e:
            print(f"Error al guardar el archivo de tipo {tipo}: {e}")
            return None
        
    def guardar_archivo_general(self, archivo, id, tipo, ruta, carpeta_main):
        if not archivo:
            print(f"No se encontró el archivo de tipo {tipo}.")
            return None

        try:
            # Extraer la extensión del archivo
            extension = archivo.name.split('.')[-1]
            filename = f'{tipo}_{carpeta_main}_{id}.{extension}'

            ruta_base = os.path.join(settings.MEDIA_ROOT, ruta, carpeta_main, f"{tipo}_{id}")
            if not os.path.exists(ruta_base):
                os.makedirs(ruta_base)
                print(f"Carpeta creada: '{ruta_base}'")
            else:
                print(f"La carpeta ya existe: '{ruta_base}'")

            # Ruta completa del archivo
            full_path = os.path.join(ruta_base, filename)

            # Guardar el archivo en el directorio especificado
            with open(full_path, 'wb') as destination:
                for chunk in archivo.chunks():
                    destination.write(chunk)

            print(f"Archivo de tipo {tipo} guardado en: '{full_path}'")
            return full_path
        except Exception as e:
            print(f"Error al guardar el archivo de tipo {tipo}: {e}")
            return None

class CursosTblViewSet(viewsets.ModelViewSet):
    queryset = Cursos.objects.all()
    serializer_class = CursosSerializer
    permission_classes = [IsAuthenticated]  # Por ejemplo, requiere autenticación para todas las acciones

    def list(self, request, *args, **kwargs):
        # Sobrescribimos el método list para personalizar la respuesta
        queryset = self.get_queryset()
        querysetado = self.queryset.filter(status=True)
        serializer = self.get_serializer(querysetado, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)

    def create(self, request, *args, **kwargs):
            datos_totales = request.data
            print("datos_totales.......................",datos_totales)
            img_logo = datos_totales.pop('img_logo', None)
            img_fondo = datos_totales.pop('img_fondo', None)
            vid_trailer = datos_totales.pop('vid_trailer', None)
            # Handle file uploads correctly
            img_logo = request.FILES.get('img_logo')
            img_fondo = request.FILES.get('img_fondo')
            vid_trailer = request.FILES.get('vid_trailer')
            print("datos_totales.......................",datos_totales)
            # Remove the file fields  
            lista_alumnos = json.loads(datos_totales.get('alumno'))  # Suponiendo que los datos de usuarios están bajo una clave llamada 'usuarios'
            # datos_totales.pop('alumno', None)
            print("lista_alumnos.......................",lista_alumnos)
            serializer = self.get_serializer(data=datos_totales)
            serializer.is_valid(raise_exception=True)
            curso = serializer.save()

            # Crear instancias de CursoUsuario para cada usuario y guardarlas en la base de datos
            for usuario in lista_alumnos:
                alum_usuario_id = usuario
                user_data = User.objects.filter(id = alum_usuario_id).first()
                # usuario_perfil = get_object_or_404(PerfilUsuario, id=user_data.id)
                curso_usuario = CursoUsuario.objects.create(usuarios=user_data, cursos=curso,status=True)
                curso_usuario.save()

            print("img logo .............",img_logo)
            print("img fondo ..........................",img_fondo)
            print("vid_trailer .................................",vid_trailer)
            id_curso = curso.id
            # Guardar archivos relacionados al curso
            file_handler = FileHandler()

            ruta_app = "/usr/src/app/"
            ruta_local = "http://localhost:8000/"

            print("logo_path .............",img_logo)
            print("fondo_path..........................",img_fondo)
            print("trailer_path .................................",vid_trailer)


            if img_logo:
                # logo_path = file_handler.guardar_archivo(img_logo, id_curso,"logo","curso")
                logo_path = file_handler.guardar_archivo_general(img_logo, id_curso, "logo_curso","cursos","logo")

                img_logo_nueva = logo_path.replace(ruta_app, ruta_local)
                curso.img_logo = img_logo_nueva if img_logo_nueva else curso.img_logo

            if img_fondo:
                # fondo_path = file_handler.guardar_archivo(img_fondo, id_curso,"fondo","curso")
                fondo_path = file_handler.guardar_archivo_general(img_fondo, id_curso, "fondo_curso","cursos","fondo")

                img_fondo_nueva = fondo_path.replace(ruta_app, ruta_local)
                curso.img_fondo = img_fondo_nueva if img_fondo_nueva else curso.img_fondo

            if vid_trailer:
                # trailer_path = file_handler.guardar_archivo(vid_trailer, id_curso,"trailer","curso")
                trailer_path = file_handler.guardar_archivo_general(vid_trailer, id_curso, "trailer_curso","cursos","trailer")

                vid_trailer_nueva = trailer_path.replace(ruta_app, ruta_local)
                curso.vid_trailer = vid_trailer_nueva if vid_trailer_nueva else curso.vid_trailer

            # actualizar el objeto curso con las rutas de los archivos
            curso.save()
            headers = self.get_success_headers(serializer.data)
            return Response(serializer.data, status=status.HTTP_201_CREATED, headers=headers)

        
    def update(self, request, *args, **kwargs):
        partial = kwargs.pop('partial', False)
        instance = self.get_object()
        datos_totales = request.data
        print("datos_totales.......................",datos_totales)

        img_logo = request.FILES.get('img_logo', None)
        if img_logo:
                datos_totales.pop('img_logo', None)

        img_fondo = request.FILES.get('img_fondo', None)
        if img_fondo:
                datos_totales.pop('img_fondo', None)

        vid_trailer = request.FILES.get('vid_trailer', None)
        if vid_trailer:
                datos_totales.pop('vid_trailer', None)


        lista_alumnos = json.loads(datos_totales.get('alumno'))  # Suponiendo que los datos de usuarios están bajo una clave llamada 'usuarios'
        # datos_totales.pop('alumno', None)
        print("lista_alumnos.......................",lista_alumnos)
        

        print("datos_totales.........",datos_totales)
        serializer = self.get_serializer(instance, data=datos_totales, partial=partial)
        serializer.is_valid(raise_exception=True)
        curso = serializer.save()
        id_curso = curso.id

        CursoUsuario.objects.filter(cursos = id_curso).update(status=False)

        for usuario in lista_alumnos:
            alum_usuario_id = usuario
            print("alum_usuario_id  ...--........",alum_usuario_id)

            # usuario_perfil = get_object_or_404(PerfilUsuario, id=user_data.id)
            # curso_usuario = CursoUsuario.objects.update(usuarios=user_data, cursos=curso)
            user_data = User.objects.filter(id=alum_usuario_id).first()
            print("user_data....update............................................",user_data)
            curso_usuario, created = CursoUsuario.objects.update_or_create(
                usuarios=user_data,
                cursos=curso,
                defaults={ 'status': True}
            )
            # , 'estado_examen_curso': 1 ,'est_avance':1
            print("aquii___")

            # Ahora puedes devolver una respuesta adecuada
            if created:
                message = 'CursoUsuario creado exitosamente.'
                print("messagei___",message)

            else:
                message = 'CursoUsuario actualizado exitosamente.'
                print("messagei___",message)



        # Guardar archivos relacionados al curso
        file_handler = FileHandler()

        ruta_app = "/usr/src/app/"
        ruta_local = "http://localhost:8000/"

        print("logo_path .............",img_logo)
        print("fondo_path..........................",img_fondo)
        print("trailer_path .................................",vid_trailer)


        if img_logo:
            logo_path = file_handler.guardar_archivo(img_logo, id_curso, "logo","curso")
            img_logo_nueva = logo_path.replace(ruta_app, ruta_local)
            curso.img_logo = img_logo_nueva if img_logo_nueva else curso.img_logo

        if img_fondo:
            fondo_path = file_handler.guardar_archivo(img_fondo, id_curso, "fondo","curso")
            img_fondo_nueva = fondo_path.replace(ruta_app, ruta_local)
            curso.img_fondo = img_fondo_nueva if img_fondo_nueva else curso.img_fondo

        if vid_trailer:
            trailer_path = file_handler.guardar_archivo(vid_trailer, id_curso, "trailer","curso")
            vid_trailer_nueva = trailer_path.replace(ruta_app, ruta_local)
            curso.vid_trailer = vid_trailer_nueva if vid_trailer_nueva else curso.vid_trailer

        # actualizar el objeto curso con las rutas de los archivos
        curso.save()
        

        headers = self.get_success_headers(serializer.data)
        return Response(serializer.data, status=status.HTTP_200_OK, headers=headers)

class ContenidoTblViewSet(viewsets.ModelViewSet):
    queryset = Contenidos.objects.all()
    serializer_class = ContenidoSerializer
    permission_classes = [IsAuthenticated]

    def list(self, request, *args, **kwargs):
        # Sobrescribimos el método list para personalizar la respuesta
        queryset = self.get_queryset()
        querysetado = self.queryset.filter(status=True)
        serializer = self.get_serializer(querysetado, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)

    def list_by_curso(self, request, pk=None):
        contenidos = Contenidos.objects.filter(curso=pk,status=True)
        serializer = self.get_serializer(contenidos, many=True)
        return Response(serializer.data)

    def create(self, request, *args, **kwargs):
        datos_totales = request.data

        print("datos_totales dos: ", datos_totales)
        serializer = self.get_serializer(data=datos_totales)
        serializer.is_valid(raise_exception=True)
        contenido = serializer.save()
        id_contenido = contenido.id
        contenido.save()

        headers = self.get_success_headers(serializer.data)
        return Response(serializer.data, status=status.HTTP_201_CREATED, headers=headers)

    def update(self, request, *args, **kwargs):
        partial = kwargs.pop('partial', False)
        instance = self.get_object()
        datos_totales = request.data
        print("datos_totales.........",datos_totales)
        serializer = self.get_serializer(instance, data=datos_totales, partial=partial)
        serializer.is_valid(raise_exception=True)
        contenido = serializer.save()
        id_contenido = contenido.id
        contenido.save()

        headers = self.get_success_headers(serializer.data)
        return Response(serializer.data, status=status.HTTP_200_OK, headers=headers)
    
    def retrieve(self, request, pk=None):
        contenido = get_object_or_404(Contenidos, pk=pk)
        serializer = self.get_serializer(contenido)
        return Response(serializer.data)
    
    @action(detail=False, methods=['delete'])
    def delete(self, request, pk, format=None):
        print("pk..........................", pk)
        try:
            self.queryset.filter(id=pk).update(status=False)
            return Response(status=status.HTTP_204_NO_CONTENT)
        except Contenidos.DoesNotExist:
            return Response({'error': 'Contenidos not found'}, status=status.HTTP_404_NOT_FOUND)
        
class TemasTblViewSet(viewsets.ModelViewSet):
    queryset = Temas.objects.all()
    serializer_class = TemasSerializer
    permission_classes = [IsAuthenticated]

    def list(self, request, *args, **kwargs):
        # Sobrescribimos el método list para personalizar la respuesta
        queryset = self.get_queryset()
        querysetado = self.queryset.filter(status=True)
        serializer = self.get_serializer(querysetado, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)

    def list_by_type(self, request, id, type=None):
        if type == 'contenido':
            archivos = self.queryset.filter(contenido=id , status=True )
        elif type == 'id':
            archivos = self.queryset.filter(id=id)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)

    def create(self, request, *args, **kwargs):
        datos_totales = request.data
        print("datos_totales:", datos_totales)

        img_fondo = request.FILES.get('img_fondo', None)
        video_tema = request.FILES.get('video_tema', None)
        multiple_files = request.FILES.getlist('multiple_files')

        if video_tema:
            datos_totales.pop('video_tema', None)

        datos_totales.pop('multiple_files', None)
        datos_totales.pop('img_fondo', None)
        print("datos_totales dos: ", datos_totales)

        file_handler = FileHandler()

        serializer = self.get_serializer(data=datos_totales)
        serializer.is_valid(raise_exception=True)
        temas = serializer.save()

        id_temas = temas.id

        ruta_app = "/usr/src/app/"
        ruta_local = "http://localhost:8000/"

        # Guardar archivos múltiples
        for archivo in multiple_files:
            print("archivo ....................",archivo)
            # archivo_path = file_handler.guardar_archivo_temas(archivo, id_temas, tipo="archivos_temas")
            archivo_path = file_handler.guardar_archivo_general(archivo, id_temas, "archivos_temas","cursos/contenido/tema","archivos_temas")

            archivo_nueva = archivo_path.replace(ruta_app, ruta_local)
            ArchivoTemas.objects.create(archivo=archivo_nueva, temas=temas)

        # Guardar imagen de fondo
        if img_fondo:
            # fondo_path = file_handler.guardar_archivo_temas(img_fondo, id_temas, tipo="fondo_temas")
            fondo_path = file_handler.guardar_archivo_general(img_fondo, id_temas, "fondo_temas","cursos/contenido/tema","fondo_temas")

            img_fondo_nueva = fondo_path.replace(ruta_app, ruta_local)
            temas.img_fondo = img_fondo_nueva

        # Guardar video contenido
        if video_tema:
            # video_path = file_handler.guardar_archivo_temas(video_tema, id_temas, tipo="video_tema")
            video_path = file_handler.guardar_archivo_general(video_tema, id_temas, "video_tema","cursos/contenido/tema","video_tema")

            video_tema_nueva = video_path.replace(ruta_app, ruta_local)
            temas.video_tema = video_tema_nueva

        temas.save()

        headers = self.get_success_headers(serializer.data)
        return Response(serializer.data, status=status.HTTP_201_CREATED, headers=headers)

    def update(self, request, *args, **kwargs):
        partial = kwargs.pop('partial', False)
        instance = self.get_object()
        datos_totales = request.data
        print("datos_totales...........................",datos_totales)

        img_fondo = request.FILES.get('img_fondo', None)
        if img_fondo:
            datos_totales.pop('img_fondo', None)

        video_tema = request.FILES.get('video_tema', None)
        if video_tema:
            datos_totales.pop('video_tema', None)

        multiple_files = request.FILES.getlist('multiple_files')
        if multiple_files:
            datos_totales.pop('multiple_files')

        file_handler = FileHandler()

        serializer = self.get_serializer(instance, data=datos_totales, partial=partial)
        serializer.is_valid(raise_exception=True)
        temas = serializer.save()

        id_temas = temas.id

        ruta_app = "/usr/src/app/"
        ruta_local = "http://localhost:8000/"

        # Guardar archivos múltiples
        if multiple_files != '':  
            # Remove existing multiple files if any and save the new ones
            # ArchivoContenido.objects.filter(temas=temas).delete()
            for archivo in multiple_files:
                # archivo_path = file_handler.guardar_archivo_temas(archivo, id_temas, tipo="archivos_temas")
                archivo_path = file_handler.guardar_archivo_general(archivo, id_temas, "archivos_temas","cursos/contenido/tema","archivos_temas")

                archivo_nueva = archivo_path.replace(ruta_app, ruta_local)
                ArchivoTemas.objects.create(archivo=archivo_nueva, temas=temas)

        # Guardar imagen de fondo
        if img_fondo:
            fondo_path = file_handler.guardar_archivo_temas(img_fondo, id_temas, tipo="fondo_temas")
            fondo_path = file_handler.guardar_archivo_general(img_fondo, id_temas, "fondo_temas","cursos/contenido/tema","fondo_temas")
            img_fondo_nueva = fondo_path.replace(ruta_app, ruta_local)
            temas.img_fondo = img_fondo_nueva

        # Guardar video temas
        if video_tema:
            # video_path = file_handler.guardar_archivo_temas(video_tema, id_temas, tipo="video_tema")
            video_path = file_handler.guardar_archivo_general(video_tema, id_temas, "video_tema","cursos/contenido/tema","video_tema")

            video_tema_nueva = video_path.replace(ruta_app, ruta_local)
            temas.video_tema = video_tema_nueva

        temas.save()

        headers = self.get_success_headers(serializer.data)
        return Response(serializer.data, status=status.HTTP_200_OK, headers=headers)


    # @action(detail=False, methods=['delete'])
    def delete(self, request, pk, format=None):
        print("pk..........................", pk)
        try:
            # archivo_temas = self.get_queryset().get(id=id_archivo_tema)
            # archivo_temas.delete()
            # Actualiza el objeto en lugar de eliminarlo
            self.queryset.filter(id=pk).update(status=False)
            listas_archivos=ArchivoTemas.objects.filter(temas=pk)
            print("listas_archivos.............................................",listas_archivos)
            return Response(status=status.HTTP_204_NO_CONTENT)
        except Temas.DoesNotExist:
            return Response({'error': 'Temas not found'}, status=status.HTTP_404_NOT_FOUND)
        
class ArchivoTemasTblViewSet(viewsets.ModelViewSet):
    queryset = ArchivoTemas.objects.all()
    serializer_class = ArchivoTemasSerializer
    permission_classes = [IsAuthenticated]

    def destroy(self, request, *args, **kwargs):
        id_archivo_tema = kwargs.get('id_archivo_tema')
        print("id_archivo_tema..........................",id_archivo_tema)
        try:
            archivo_temas = self.get_queryset().get(id=id_archivo_tema)
            archivo = archivo_temas.archivo
            ruta_local = "http://localhost:8000/media"
            path_main = archivo.replace(ruta_local,'')
            path_root = os.path.join(settings.MEDIA_ROOT)
            path_main=  path_root + path_main
            print("path_main ************************",path_main)
            # Delete the file from the filesystem
            if os.path.isfile(path_main):
                os.remove(path_main)

            archivo_temas.delete()
            return Response(status=status.HTTP_204_NO_CONTENT)
        except ArchivoTemas.DoesNotExist:

            return Response({'error': 'ArchivoTemas not found'}, status=status.HTTP_404_NOT_FOUND)
        
    def list_by_type(self, request, id, type=None):
        if type == 'curso':
            archivos = self.queryset.filter(contenido=id)
        # elif type == 'quiz':
        #     archivos = self.queryset.filter(quiz_id=id)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)
    
class QuizTblViewSet(viewsets.ModelViewSet):
    queryset = Quiz.objects.all()
    serializer_class = QuizSerializer
    permission_classes = [IsAuthenticated]


    @action(detail=False, methods=['get'])
        # def list(self, request):
        #     # Custom logic for listing all quizzes
        #     queryset = self.get_queryset()
        #     serializer = self.get_serializer(queryset, many=True)
        #     return Response(serializer.data)
    def list(self, request, *args, **kwargs):
        # Sobrescribimos el método list para personalizar la respuesta
        queryset = self.get_queryset()
        querysetado = self.queryset.filter(status=True)
        serializer = self.get_serializer(querysetado, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)


    @action(detail=False, methods=['get'])
    def list_id(self, request):
        # Custom logic for listing IDs
        data = Quiz.objects.values_list('id', flat=True)
        return Response(data)


    
    @action(detail=False, methods=['get'])
    def list_by_type(self, request, id, type=None):
        if type == 'contenido':
            archivos = self.queryset.filter(contenido=id,status=True)
        elif type == 'id':
            archivos = self.queryset.filter(id=id)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)
    
    @action(detail=False, methods=['post'])
    def create(self, request, *args, **kwargs):
        # datos_totales = request.data.copy()  # Make a mutable copy of the request data
        datos_totales = request.data  # Make a mutable copy of the request data
        print("datos_totales.............................",datos_totales)
        imagen_quiz = request.FILES.get('imagen_quiz', None)
        print("imagen_quiz.........",imagen_quiz)
        if imagen_quiz:
            datos_totales.pop('imagen_quiz', None)

        # Custom logic for creating a quiz
        serializer = self.get_serializer(data=datos_totales)
        serializer.is_valid(raise_exception=True)
        quiz = serializer.save()

        if imagen_quiz:
            file_handler = FileHandler()
            ruta_app = "/usr/src/app/"
            ruta_local = "http://localhost:8000/"
            # imagen_path = file_handler.guardar_archivo_quiz(imagen_quiz, quiz.id, "imagen_quiz")
            imagen_path = file_handler.guardar_archivo_general(imagen_quiz, quiz.id, "imagen_quiz","cursos/contenido/quiz","imagen_quiz")

            imagen_nueva = imagen_path.replace(ruta_app, ruta_local)
            quiz.imagen_quiz = imagen_nueva
            quiz.save()

        headers = self.get_success_headers(serializer.data)
        return Response(serializer.data, status=status.HTTP_201_CREATED, headers=headers)
    
    @action(detail=False, methods=['put'])
    def update(self, request, *args, **kwargs):
        partial = kwargs.pop('partial', False)
        datos_totales = request.data  # Make a mutable copy of the request data

        imagen_quiz = request.FILES.get('imagen_quiz', None)
        print("imagen_quiz.........",imagen_quiz)
        if imagen_quiz:
            datos_totales.pop('imagen_quiz', None)


        # Custom logic for updating a quiz
        instance = self.get_object()
        serializer = self.get_serializer(instance, data=request.data, partial=partial)
        serializer.is_valid(raise_exception=True)
        quiz = serializer.save()

        self.perform_update(serializer)

        if imagen_quiz:
            file_handler = FileHandler()
            ruta_app = "/usr/src/app/"
            ruta_local = "http://localhost:8000/"
            # imagen_path = file_handler.guardar_archivo_quiz(imagen_quiz, quiz.id, "imagen_quiz")
            imagen_path = file_handler.guardar_archivo_general(imagen_quiz, quiz.id, "imagen_quiz","cursos/contenido/quiz","imagen_quiz")

            imagen_nueva = imagen_path.replace(ruta_app, ruta_local)
            quiz.imagen_quiz = imagen_nueva
            quiz.save()


        return Response(serializer.data)

    @action(detail=False, methods=['delete'])
    # def destroy(self, request):
    #     # Custom logic for deleting a quiz
    #     instance = self.get_object()
    #     self.perform_destroy(instance)
    #     return Response(status=status.HTTP_204_NO_CONTENT)

    def delete(self, request, pk, format=None):
        print("pk..........................", pk)
        try:
            # archivo_temas = self.get_queryset().get(id=id_archivo_tema)
            # archivo_temas.delete()
            # Actualiza el objeto en lugar de eliminarlo
            self.queryset.filter(id=pk).update(status=False)
            return Response(status=status.HTTP_204_NO_CONTENT)
        except Quiz.DoesNotExist:
            return Response({'error': 'Quiz not found'}, status=status.HTTP_404_NOT_FOUND) 
  
class PreguntasTblViewSet(viewsets.ModelViewSet):
    queryset = PreguntaQuiz.objects.all()
    serializer_class = PreguntaQuizSerializer
    permission_classes = [IsAuthenticated]

    @action(detail=False, methods=['get'])
    def list_id(self, request):
        data = PreguntaQuiz.objects.values_list('id', flat=True)
        return Response(data)

    @action(detail=False, methods=['get'])
    def list(self, request, *args, **kwargs):
        queryset = self.get_queryset()
        querysetado = self.queryset.filter(status=True)
        serializer = self.get_serializer(querysetado, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)

    @action(detail=False, methods=['get'])
    def list_by_type(self, request, id, type=None):
        if type == 'quiz':
            archivos = self.queryset.filter(quiz=id,status=True)
        # elif type == 'quiz':
        #     archivos = self.queryset.filter(quiz_id=id)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)
    
    @action(detail=False, methods=['post'])
    def create(self, request, *args, **kwargs):
        # datos_totales = request.data.copy()  # Make a mutable copy of the request data
        datos_totales = request.data  # Make a mutable copy of the request data
        # datos_totales = request.data.copy()
        print("datos_totales...............",datos_totales)

        # Convierte los datos totales en una QueryDict
        # query_dict = QueryDict(mutable=True)
        # for key, value in datos_totales.items():
        #     if key != 'opciones_correctas':
        #         query_dict.update({key: value})
                                
        # opciones_correctas_str = request.data.get('opciones_correctas', None)
        # opciones_data = json.loads(opciones_correctas_str)    

        imagen_pregunta = request.FILES.get('imagen_pregunta', None)


        if imagen_pregunta:
            datos_totales.pop('imagen_pregunta', None)

        serializer = self.get_serializer(data=datos_totales)
        serializer.is_valid(raise_exception=True)
        pregunta = serializer.save()

        if imagen_pregunta:
            file_handler = FileHandler()
            ruta_app = "/usr/src/app/"
            ruta_local = "http://localhost:8000/"
            # imagen_path = file_handler.guardar_archivo_preguntas(imagen_pregunta, pregunta.id, tipo="imagen_pregunta")
            imagen_path = file_handler.guardar_archivo_general(imagen_pregunta, pregunta.id, "imagen_pregunta","cursos/contenido/quiz/pregunta","imagen_pregunta")

            imagen_nueva = imagen_path.replace(ruta_app, ruta_local)
            pregunta.imagen_pregunta = imagen_nueva
            pregunta.save()

        headers = self.get_success_headers(serializer.data)
        return Response(serializer.data, status=status.HTTP_201_CREATED, headers=headers)
    

    @action(detail=False, methods=['put'])
    def update(self, request, *args, **kwargs):
        partial = kwargs.pop('partial', False)
        datos_totales = request.data  # Make a mutable copy of the request data
        imagen_pregunta = request.FILES.get('imagen_pregunta', None)
        print("imagen_pregunta.........",imagen_pregunta)
        if imagen_pregunta:
            datos_totales.pop('imagen_pregunta', None)

        # Custom logic for updating a pregunta
        instance = self.get_object()
        serializer = self.get_serializer(instance, data=request.data, partial=partial)
        serializer.is_valid(raise_exception=True)
        pregunta = serializer.save()

        self.perform_update(serializer)

        if imagen_pregunta:
            file_handler = FileHandler()
            ruta_app = "/usr/src/app/"
            ruta_local = "http://localhost:8000/"
            # imagen_path = file_handler.guardar_archivo_preguntas(imagen_pregunta, pregunta.id, tipo="imagen_pregunta")
            imagen_path = file_handler.guardar_archivo_general(imagen_pregunta, pregunta.id, "imagen_pregunta","cursos/contenido/quiz/pregunta","imagen_pregunta")

            imagen_nueva = imagen_path.replace(ruta_app, ruta_local)
            pregunta.imagen_pregunta = imagen_nueva
            pregunta.save()

        return Response(serializer.data)


    @action(detail=False, methods=['delete'])
    def delete(self, request, pk, format=None):
        print("pk..........................", pk)
        try:
            self.queryset.filter(id=pk).update(status=False)
            return Response(status=status.HTTP_204_NO_CONTENT)
        except PreguntaQuiz.DoesNotExist:
            return Response({'error': 'PreguntaQuiz not found'}, status=status.HTTP_404_NOT_FOUND)
        
class ExamenCursoViewSet(viewsets.ModelViewSet):
    queryset = ExamenCurso.objects.all()
    serializer_class = ExamenCursoSerializer
    permission_classes = [IsAuthenticated]

    @action(detail=False, methods=['get'])
    def list(self, request, *args, **kwargs):
        queryset = self.get_queryset()
        querysetado = self.queryset.filter(status=True)
        serializer = self.get_serializer(querysetado, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)

    @action(detail=False, methods=['get'])
    def list_id(self, request):
        # Custom logic for listing IDs
        data = ExamenCurso.objects.values_list('id', flat=True)
        return Response(data)
    
    @action(detail=False, methods=['get'])
    def list_by_type(self, request, id, type=None):
        if type == 'curso':
            archivos = self.queryset.filter(curso=id,status=True)
        elif type == 'id':
            archivos = self.queryset.filter(id=id)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)

    @action(detail=False, methods=['post'])
    def create(self, request, *args, **kwargs):
        datos_totales = request.data  
        imagen_examen = request.FILES.get('imagen_examen', None)
        if imagen_examen:
            datos_totales.pop('imagen_examen', None)

        serializer = self.get_serializer(data=datos_totales)
        serializer.is_valid(raise_exception=True)
        examen = serializer.save()

        if imagen_examen:
            file_handler = FileHandler()
            ruta_app = "/usr/src/app/"
            ruta_local = "http://localhost:8000/"
            imagen_path = file_handler.guardar_archivo_general(imagen_examen, examen.id, "imagen_examen","cursos","examen")
            imagen_nueva = imagen_path.replace(ruta_app, ruta_local)
            examen.imagen_examen = imagen_nueva
            examen.save()

        headers = self.get_success_headers(serializer.data)
        return Response(serializer.data, status=status.HTTP_201_CREATED, headers=headers)
    
    @action(detail=False, methods=['put'])
    def update(self, request, *args, **kwargs):
        partial = kwargs.pop('partial', False)
        datos_totales = request.data  # Make a mutable copy of the request data
        imagen_examen = request.FILES.get('imagen_examen', None)
        if imagen_examen:
            datos_totales.pop('imagen_examen', None)


        # Custom logic for updating a examen
        instance = self.get_object()
        serializer = self.get_serializer(instance, data=request.data, partial=partial)
        serializer.is_valid(raise_exception=True)
        examen = serializer.save()

        self.perform_update(serializer)

        if imagen_examen:
            file_handler = FileHandler()
            ruta_app = "/usr/src/app/"
            ruta_local = "http://localhost:8000/"
            # imagen_path = file_handler.guardar_archivo_general(imagen_examen, examen.id, "imagen_examen")
            imagen_path = file_handler.guardar_archivo_general(imagen_examen, examen.id, "imagen_examen","cursos","examen")

            imagen_nueva = imagen_path.replace(ruta_app, ruta_local)
            examen.imagen_examen = imagen_nueva
            examen.save()


        return Response(serializer.data)

    @action(detail=False, methods=['delete'])

    def delete(self, request, pk, format=None):
        print("pk..........................", pk)
        try:
            # archivo_temas = self.get_queryset().get(id=id_archivo_tema)
            # archivo_temas.delete()
            # Actualiza el objeto en lugar de eliminarlo
            self.queryset.filter(id=pk).update(status=False)
            return Response(status=status.HTTP_204_NO_CONTENT)
        except ExamenCurso.DoesNotExist:
            return Response({'error': 'ExamenCurso not found'}, status=status.HTTP_404_NOT_FOUND)
                
class PreguntasExamenCursoViewSet(viewsets.ModelViewSet):
    queryset = PreguntaExamenCurso.objects.all()
    serializer_class = PreguntasExamenCursoSerializer
    permission_classes = [IsAuthenticated]

    @action(detail=False, methods=['get'])
    def list_id(self, request):
        data = PreguntaExamenCurso.objects.values_list('id', flat=True)
        return Response(data)

    @action(detail=False, methods=['get'])
    def list(self, request, *args, **kwargs):
        queryset = self.get_queryset()
        querysetado = self.queryset.filter(status=True)
        serializer = self.get_serializer(querysetado, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)

    @action(detail=False, methods=['get'])
    def list_by_type(self, request, id, type=None):
        if type == 'examen_curso':
            archivos = self.queryset.filter(examen_curso=id,status=True)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)
    
    @action(detail=False, methods=['post'])
    def create(self, request, *args, **kwargs):
        datos_totales = request.data
        print("datos_totales...............",datos_totales)
        imagen_pregunta = request.FILES.get('imagen_pregunta', None)

        if imagen_pregunta:
            datos_totales.pop('imagen_pregunta', None)

        serializer = self.get_serializer(data=datos_totales)
        serializer.is_valid(raise_exception=True)
        pregunta = serializer.save()

        if imagen_pregunta:
            file_handler = FileHandler()
            ruta_app = "/usr/src/app/"
            ruta_local = "http://localhost:8000/"
            imagen_path = file_handler.guardar_archivo_general(imagen_pregunta, pregunta.id, "imagen_pregunta_examen","cursos/examen","pregunta_examen")

            imagen_nueva = imagen_path.replace(ruta_app, ruta_local)
            pregunta.imagen_pregunta = imagen_nueva
            pregunta.save()

        headers = self.get_success_headers(serializer.data)
        return Response(serializer.data, status=status.HTTP_201_CREATED, headers=headers)
    

    @action(detail=False, methods=['put'])
    def update(self, request, *args, **kwargs):
        partial = kwargs.pop('partial', False)
        datos_totales = request.data  # Make a mutable copy of the request data
        imagen_pregunta = request.FILES.get('imagen_pregunta', None)
        print("imagen_pregunta.........",imagen_pregunta)
        if imagen_pregunta:
            datos_totales.pop('imagen_pregunta', None)

        # Custom logic for updating a pregunta
        instance = self.get_object()
        serializer = self.get_serializer(instance, data=request.data, partial=partial)
        serializer.is_valid(raise_exception=True)
        pregunta = serializer.save()

        self.perform_update(serializer)

        if imagen_pregunta:
            file_handler = FileHandler()
            ruta_app = "/usr/src/app/"
            ruta_local = "http://localhost:8000/"
            imagen_path = file_handler.guardar_archivo_general(imagen_pregunta, pregunta.id, "imagen_pregunta_examen","cursos/examen","pregunta_examen")

            imagen_nueva = imagen_path.replace(ruta_app, ruta_local)
            pregunta.imagen_pregunta = imagen_nueva
            pregunta.save()

        return Response(serializer.data)


    @action(detail=False, methods=['delete'])
    def delete(self, request, pk, format=None):
        print("pk..........................", pk)
        try:
            self.queryset.filter(id=pk).update(status=False)
            return Response(status=status.HTTP_204_NO_CONTENT)
        except PreguntaExamenCurso.DoesNotExist:
            return Response({'error': 'PreguntaExamenCurso not found'}, status=status.HTTP_404_NOT_FOUND)
        
class CursoUsuarioTblViewSet(viewsets.ModelViewSet):
    queryset = CursoUsuario.objects.all()
    serializer_class = CursoUsuarioMainSerializer
    permission_classes = [IsAuthenticated]

    @action(detail=False, methods=['get'])
    def list_id(self, request):
        data = CursoUsuario.objects.values_list('id', flat=True)
        return Response(data)

    @action(detail=False, methods=['get'])
    def list(self, request, *args, **kwargs):
        queryset = self.get_queryset()
        querysetado = self.queryset.filter(status=True)
        serializer = self.get_serializer(querysetado, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)

    @action(detail=False, methods=['get'])
    def list_by_type(self, request, id, type=None):
        if type == 'usuario':
            archivos = self.queryset.filter(usuarios=id,status=True)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)
    
    @action(detail=False, methods=['post'])
    def registrar_acciones(self, request, *args, **kwargs):
        datos_totales = request.data  
        print("datos... usuario curso examen.........................",datos_totales)

        tipo= datos_totales['accion']
        if tipo == 'registro_acceso_examen':
            serializer = self.get_serializer(data=datos_totales)
            serializer.is_valid(raise_exception=True)
            curso_usu = serializer.save()
            curso_usu.estado_examen_curso = 2
            curso_usu.save()
            headers = self.get_success_headers(serializer.data)
            
        elif tipo == 'registro_termino_examen':
            # archivos = self.queryset.filter(id=id)
            print("............registro_termino_examen..........")
            serializer = self.get_serializer(data=datos_totales)
            serializer.is_valid(raise_exception=True)
            curso_usu = serializer.save()
            curso_usu.estado_examen_curso = 3
            curso_usu.est_avance=3
            curso_usu.save()
            headers = self.get_success_headers(serializer.data)

        elif tipo == 'otra_accion':
            # archivos = self.queryset.filter(id=id)
            print("otra accion...............")

        else:
            return Response({"detail": "Invalid type"}, status=400)
  
        return Response('registrado', status=status.HTTP_201_CREATED)

    @action(detail=False, methods=['put'])
    def actualizar_acciones(self, request, *args, **kwargs):
        partial = kwargs.pop('partial', False)
        datos_totales = request.data.copy()
        # request.data 
        instance = self.get_object()
 
        # print("accion...............",datos_totales['accion'])
        print("partial...............",partial)
        print("instance...............",instance)


        tipo= datos_totales['accion']
        print("tipo...............",tipo)
        datos_totales.pop('accion', None)

        if tipo == 'registro_acceso_examen':
            print("tipo...............",datos_totales)
            datos_totales['usuarios'] = instance.usuarios.id
            datos_totales['estado_examen_curso'] = 2
            # datos_totales['cursos'] = instance.cursos.id
            print("datos_totales...............",datos_totales)
            serializer = self.get_serializer(instance, data=datos_totales, partial=True)
            serializer.is_valid(raise_exception=True)
            curso_usu = serializer.save()

            # curso_usu.estado_examen_curso = 2
            # curso_usu.save()
            self.perform_update(serializer)

        elif tipo == 'registro_termino_examen':
            # archivos = self.queryset.filter(id=id)
            print("tipo...............",datos_totales)
            datos_totales['usuarios'] = instance.usuarios.id
            datos_totales['estado_examen_curso'] = 3
            datos_totales['est_avance'] = 3
            # datos_totales['cursos'] = instance.cursos.id
            print("datos_totales...............",datos_totales)
            serializer = self.get_serializer(instance, data=datos_totales, partial=True)
            serializer.is_valid(raise_exception=True)
            curso_usu = serializer.save()

            # curso_usu.estado_examen_curso = 2
            # curso_usu.save()
            self.perform_update(serializer)

        elif tipo == 'actualizar_estrellas':
            # archivos = self.queryset.filter(id=id)
            print("tipo...............",datos_totales)
            # datos_totales['feedback_puntaje'] = instance.usuarios.id
            # datos_totales['cursos'] = instance.cursos.id
            print("datos_totales...............",datos_totales)
            serializer = self.get_serializer(instance, data=datos_totales, partial=True)
            serializer.is_valid(raise_exception=True)
            curso_usu = serializer.save()

            # curso_usu.estado_examen_curso = 2
            # curso_usu.save()
            self.perform_update(serializer)

        elif tipo == 'otra_accion':
            # archivos = self.queryset.filter(id=id)
            print("otra accion...............")

        else:

            print("else...............")

        return Response('Actualizado')

    def retrieve(self, request, pk=None,usu=None):
        cursos_usu_dat = get_object_or_404(CursoUsuario, pk=pk)
        serializer = self.get_serializer(cursos_usu_dat)

        # print("serializer.data",serializer.data)
        return Response(serializer.data)



class NotasExamenViewSet(viewsets.ModelViewSet):
    queryset = NotaExamenCurso.objects.all()
    serializer_class = NotaExamenCursoSerializer
    permission_classes = [IsAuthenticated]

    @action(detail=False, methods=['get'])
    def list_id(self, request):
        data = NotaExamenCurso.objects.values_list('id', flat=True)
        return Response(data)

    @action(detail=False, methods=['get'])
    def list(self, request, *args, **kwargs):
        queryset = self.get_queryset()
        querysetado = self.queryset.filter(status=True)
        serializer = self.get_serializer(querysetado, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)

    @action(detail=False, methods=['get'])
    def list_by_type(self, request, id, type=None):
        if type == 'examen_usuario':
            archivos = self.queryset.filter(curso_usuario=id,status=True)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)
    
    @action(detail=False, methods=['post'])
    def create(self, request, *args, **kwargs):

        datos_totales = request.data

        curso_usuario_id = datos_totales.get('curso_usuario')
        examen_curso_id = datos_totales.get('examen_curso')
        usuario_reg = datos_totales.get('usuario_reg')

        respuestas_json = datos_totales.get('respuestas', [])
        respuestas = json.loads(respuestas_json)

        print("respuestas...........----.......",respuestas)

        status_texto = datos_totales.get('status')
        funcionesbasicasVistas = FuncionesBasicasVista()
        status = funcionesbasicasVistas.convertir_a_booleano(status_texto)

        # Crear el registro de NotaExamenCurso
        fecha_creado=timezone.now()

        nota_examen, creado =  NotaExamenCurso.objects.get_or_create(
            curso_usuario_id=curso_usuario_id,
            examen_curso_id=examen_curso_id,
            usuario_reg=usuario_reg,
            status=status,
        )
        nota_total = 0
        print("nota_examen..............................",nota_examen)
        print("nota_total.ini............",nota_total)

        for respuesta in respuestas:
            pregunta_id = respuesta.get('preguntaId')
            print("pregunta_id..................",pregunta_id)
            opciones = respuesta.get('opciones', [])
            pregunta = PreguntaExamenCurso.objects.get(id=pregunta_id)
            opciones_correctas = list(map(int, pregunta.opciones_correctas.strip('[]').split(',')))
            puntaje = funcionesbasicasVistas.calcular_puntaje(opciones,opciones_correctas)
            print("puntaje....",puntaje)
            nota_total += puntaje
            print("nota_total....",nota_total)
            indices = [index for index, valor in enumerate(opciones) if valor]
            print("indices....",indices)

            RespuestaUsuarioExamen.objects.get_or_create(
                nota_examen=nota_examen,
                pregunta_examen=pregunta,
                opcion_marcada=indices,
                status=status,
                usuario_reg=usuario_reg,
            )

        print("nota_total........................",nota_total)

        nota_examen.nota = nota_total
        nota_examen.creado = fecha_creado
        nota_examen.save()

        serializer = NotaExamenCursoSerializer(nota_examen)
        return Response(serializer.data)

    @action(detail=False, methods=['put'])
    def update(self, request, *args, **kwargs):
        partial = kwargs.pop('partial', False)
        datos_totales = request.data  
        print("datos_totales...............",datos_totales)

        # Custom logic for updating a pregunta
        instance = self.get_object()
        serializer = self.get_serializer(instance, data=request.data, partial=partial)
        serializer.is_valid(raise_exception=True)
        pregunta = serializer.save()

        self.perform_update(serializer)

        pregunta.imagen_pregunta = 'imagen_nueva'
        pregunta.save()

        return Response(serializer.data)


    @action(detail=False, methods=['delete'])
    def delete(self, request, pk, format=None):
        print("pk..........................", pk)
        try:
            self.queryset.filter(id=pk).update(status=False)
            return Response(status=status.HTTP_204_NO_CONTENT)
        except NotaExamenCurso.DoesNotExist:
            return Response({'error': 'NotaExamenCurso not found'}, status=status.HTTP_404_NOT_FOUND)
        
class RespuestaUsuarioExamenViewSet(viewsets.ModelViewSet):
    queryset = RespuestaUsuarioExamen.objects.all()
    serializer_class = RespuestaUsuarioExamenSerializer



class CursoEstadisticaTblViewSet(viewsets.ModelViewSet):
    queryset = Cursos.objects.all()
    serializer_class = CursoEstadisticaSerializer
    permission_classes = [IsAuthenticated]

    @action(detail=False, methods=['get'])
    def list(self, request, *args, **kwargs):
        queryset = self.get_queryset()
        querysetado = self.queryset.filter(status=True)
        serializer = self.get_serializer(querysetado, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)
    
    @action(detail=False, methods=['get'])
    def list_id(self, request,pk):

        datas = self.queryset.filter(id=pk , status=True)
        serializer = self.get_serializer(datas, many=True)
        return Response(serializer.data)    
    
    @action(detail=False, methods=['get'])
    def list_by_type(self, request, id, type=None):
        if type == 'usuario':
            archivos = self.queryset.filter(usuarios=id,status=True)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)

###############################################################################################################################################################################
###############################################################################################################################################################################
#reporte excel

class FuncionesBasicasExcelFiltros:

    # Generar un código aleatorio
    def generar_codigo_aleatorio(self,length=8):
        caracteres = string.ascii_uppercase + string.digits
        return ''.join(random.choice(caracteres) for _ in range(length))


    def extraer_archivos(self, archivo_enviado, carpeta_destino):
        try:
            os.makedirs(carpeta_destino, exist_ok=True)
            nombre_archivo = archivo_enviado.name
            archivo_path = os.path.join(carpeta_destino, nombre_archivo)
            with open(archivo_path, 'wb+') as destination:
                for chunk in archivo_enviado.chunks():
                    destination.write(chunk)
            with zipfile.ZipFile(archivo_path, 'r') as zip_file:
                # Extract all files from the archive
                zip_file.extractall(carpeta_destino)
            # Verificar y mover archivos si se crea una subcarpeta adicional
            for root, dirs, files in os.walk(carpeta_destino):
                if root != carpeta_destino:
                    for file in files:
                        os.rename(os.path.join(root, file), os.path.join(carpeta_destino, file))
                    os.rmdir(root)
            # print(f'Archivos extraídos en la carpeta: {carpeta_destino}')
            os.remove(archivo_path)
        except Exception as e:
            print(f'Error al extraer archivos: {str(e)}')


    # # Función para convertir un archivo .xls a .xlsx, reordenar columnas, filtrar por rango de fechas y por valor de Message
    # def convert_reorder_and_filter(self,file_path, start_date, end_date):
    #     try:

    #         # Lista de valores a eliminar de la columna 'Application'
    #         applications_to_exclude = ['BDA', 'CALIDAD-RED', 'CONTROL-M', 'REDOYM', 'SEGURIDAD', 'SYM']

    #         column_order = [
    #             'Run Counter', 'Origin', 'Severity', 'Status', 'Notes', 'Job Name', 'Time',
    #             'Message', 'Member/File Name', 'Application', 'Sub Application',
    #             'Control-M Server', 'Alert ID', 'Host ID', 'Order ID', 'Changed By',
    #             'Update Time', 'Run As', 'Remedy Ticket', 'Closed From EM'
    #         ]
            
    #         # Leer el archivo .xls
    #         df = pd.read_excel(file_path, engine='xlrd')  # Usa xlrd como motor de lectura

    #         # Verificar si todas las columnas requeridas están presentes
    #         required_columns = set(column_order)
    #         existing_columns = set(df.columns)
    #         if not required_columns.issubset(existing_columns):
    #             print("entro en el if ..............")
    #             missing_columns = required_columns - existing_columns
    #             raise ValueError(f"El archivo {file_path} no contiene las siguientes columnas requeridas: {missing_columns}")
            
    #         # Reordenar las columnas según column_order
    #         df = df[column_order]

    #         # Convertir las columnas a string antes de aplicar cualquier operación con .str
    #         # for col in df.columns:
    #         #     df[col] = df[col].astype(str)

    #         # Convertir la columna 'Time' a datetime
    #         df['Time'] = pd.to_datetime(df['Time'], format='%Y-%m-%d %H:%M:%S')
            
    #         # Convertir las fechas de inicio y fin a objetos datetime
    #         start_date = pd.to_datetime(start_date)
    #         end_date = pd.to_datetime(end_date + ' 23:59:59')
            
    #         # Filtrar las filas donde la columna 'Time' está dentro del rango de fechas
    #         df_filtered = df[(df['Time'] >= start_date) & (df['Time'] <= end_date)]
            
    #         # Filtrar las filas donde el valor de 'Message' sea 'Ended not OK'
    #         df_filtered = df_filtered[df_filtered['Message'] == 'Ended not OK']
            
    #         # Filtrar las filas donde el valor de 'Application' no esté en la lista de exclusión
    #         df_filtered = df_filtered[~df_filtered['Application'].isin(applications_to_exclude)]
            
    #         # Reemplazar el texto en la columna 'Notes'
    #         df_filtered['Notes'] = df_filtered['Notes'].str.replace(r'Alerta asociada al', 'Job errado')
            
    #         return df_filtered
        
    #     except Exception as e:
    #         print(f"Error al procesar el archivo {file_path}: {str(e)}")
    #         return None



    # def convert_reorder_and_filter(self, file_path, start_date, end_date):
    #     try:
    #         applications_to_exclude = ['BDA', 'CALIDAD-RED', 'CONTROL-M', 'REDOYM', 'SEGURIDAD', 'SYM']
    #         column_order = [
    #             'Run Counter', 'Origin', 'Severity', 'Status', 'Notes', 'Job Name', 'Time',
    #             'Message', 'Member/File Name', 'Application', 'Sub Application',
    #             'Control-M Server', 'Alert ID', 'Host ID', 'Order ID', 'Changed By',
    #             'Update Time', 'Run As', 'Remedy Ticket', 'Closed From EM'
    #         ]
            
    #         df = pd.read_excel(file_path, engine='xlrd')  # Leer el archivo .xls
            
    #         # Verificar que todas las columnas requeridas están presentes
    #         required_columns = set(column_order)
    #         existing_columns = set(df.columns)
    #         if not required_columns.issubset(existing_columns):
    #             missing_columns = required_columns - existing_columns
    #             raise ValueError(f"El archivo {file_path} no contiene las siguientes columnas requeridas: {missing_columns}")
            
    #         # Reordenar las columnas
    #         df = df[column_order]

    #         print("df...............",df)
            
    #         # Rellenar NaN con vacío, si no esperas NaN
    #         df.fillna('', inplace=True)
            
    #         # Convertir solo las columnas de texto en string, pero preservando los valores correctos
    #         text_columns = ['Notes', 'Job Name', 'Message', 'Member/File Name', 'Application', 'Sub Application', 'Remedy Ticket']
    #         df[text_columns] = df[text_columns].applymap(lambda x: str(x) if pd.notna(x) else x)
            
    #         # Convertir la columna 'Time' a datetime, manejando NaN
    #         df['Time'] = pd.to_datetime(df['Time'], errors='coerce', format='%Y-%m-%d %H:%M:%S')
            
    #         # Convertir las fechas de inicio y fin a objetos datetime
    #         start_date = pd.to_datetime(start_date)
    #         end_date = pd.to_datetime(end_date + ' 23:59:59')
            
    #         # Filtrar por rango de fechas en la columna 'Time'
    #         df_filtered = df[(df['Time'] >= start_date) & (df['Time'] <= end_date)]
            
    #         # Filtrar por 'Message' y 'Application'
    #         df_filtered = df_filtered[df_filtered['Message'] == 'Ended not OK']
    #         df_filtered = df_filtered[~df_filtered['Application'].isin(applications_to_exclude)]
            
    #         # Reemplazar el texto en la columna 'Notes'
    #         df_filtered['Notes'] = df_filtered['Notes'].str.replace(r'Alerta asociada al', 'Job errado', regex=True)
            
    #         return df_filtered
        
    #     except Exception as e:
    #         print(f"Error al procesar el archivo {file_path}: {str(e)}")
    #         return None


    # def convert_reorder_and_filter(self, file_path, start_date, end_date):
    #     try:
    #         applications_to_exclude = ['BDA', 'CALIDAD-RED', 'CONTROL-M', 'REDOYM', 'SEGURIDAD', 'SYM']
    #         column_order = [
    #             'Run Counter', 'Origin', 'Severity', 'Status', 'Notes', 'Job Name', 'Time',
    #             'Message', 'Member/File Name', 'Application', 'Sub Application',
    #             'Control-M Server', 'Alert ID', 'Host ID', 'Order ID', 'Changed By',
    #             'Update Time', 'Run As', 'Remedy Ticket', 'Closed From EM'
    #         ]
            
    #         # Seleccionar el motor adecuado según la extensión del archivo
    #         if file_path.endswith('.xls'):
    #             df = pd.read_excel(file_path, engine='xlrd')
    #         elif file_path.endswith('.xlsx'):
    #             df = pd.read_excel(file_path, engine='openpyxl')
    #         else:
    #             raise ValueError("Formato de archivo no soportado")

    #         # Verificar que todas las columnas requeridas están presentes
    #         required_columns = set(column_order)
    #         existing_columns = set(df.columns)
    #         if not required_columns.issubset(existing_columns):
    #             missing_columns = required_columns - existing_columns
    #             raise ValueError(f"El archivo {file_path} no contiene las siguientes columnas requeridas: {missing_columns}")
            
    #         # Reordenar las columnas
    #         df = df[column_order]

    #         # print("df...............", df)
            
    #         # Rellenar NaN con vacío
    #         df.fillna('', inplace=True)
            
    #         # Convertir solo las columnas de texto en string
    #         text_columns = ['Notes', 'Job Name', 'Message', 'Member/File Name', 'Application', 'Sub Application', 'Remedy Ticket']
    #         df[text_columns] = df[text_columns].applymap(lambda x: str(x) if pd.notna(x) else x)
            
    #         # Convertir la columna 'Time' a datetime
    #         df['Time'] = pd.to_datetime(df['Time'], errors='coerce', format='%Y-%m-%d %H:%M:%S')
            
    #         # Convertir las fechas de inicio y fin a objetos datetime
    #         start_date = pd.to_datetime(start_date)
    #         end_date = pd.to_datetime(end_date + ' 23:59:59')
            
    #         # Filtrar por rango de fechas en la columna 'Time'
    #         df_filtered = df[(df['Time'] >= start_date) & (df['Time'] <= end_date)]
            
    #         # Filtrar por 'Message' y 'Application'
    #         df_filtered = df_filtered[df_filtered['Message'] == 'Ended not OK']
    #         df_filtered = df_filtered[~df_filtered['Application'].isin(applications_to_exclude)]
            
    #         # Reemplazar el texto en la columna 'Notes'
    #         df_filtered['Notes'] = df_filtered['Notes'].str.replace(r'Alerta asociada al', 'Job errado', regex=True)
            
    #         return df_filtered
        
    #     except Exception as e:
    #         print(f"Error al procesar el archivo {file_path}: {str(e)}")
    #         return None



    # def convert_reorder_and_filter(self, file_path, start_date, end_date):
    #     try:
    #         applications_to_exclude = ['BDA', 'CALIDAD-RED', 'CONTROL-M', 'REDOYM', 'SEGURIDAD', 'SYM']
    #         column_order = [
    #             'Run Counter', 'Origin', 'Severity', 'Status', 'Notes', 'Job Name', 'Time',
    #             'Message', 'Member/File Name', 'Application', 'Sub Application',
    #             'Control-M Server', 'Alert ID', 'Host ID', 'Order ID', 'Changed By',
    #             'Update Time', 'Run As', 'Remedy Ticket', 'Closed From EM'
    #         ]
            
    #         if file_path.endswith('.xls'):
    #             df = pd.read_excel(file_path, engine='xlrd')
    #         elif file_path.endswith('.xlsx'):
    #             df = pd.read_excel(file_path, engine='openpyxl')
    #         else:
    #             raise ValueError("Formato de archivo no soportado")
            
    #         required_columns = set(column_order)
    #         if not required_columns.issubset(df.columns):
    #             missing_columns = required_columns - set(df.columns)
    #             raise ValueError(f"Columnas faltantes: {missing_columns}")
            
    #         df = df[column_order]
    #         text_columns = df.select_dtypes(include=['object', 'string']).columns
    #         df[text_columns] = df[text_columns].fillna('')
    #         df['Time'] = pd.to_datetime(df['Time'], errors='coerce')
            
    #         if df['Time'].isna().all():
    #             raise ValueError("La columna 'Time' no contiene valores válidos.")
            
    #         start_date = pd.to_datetime(start_date)
    #         end_date = pd.to_datetime(end_date + ' 23:59:59')
    #         df_filtered = df[(df['Time'] >= start_date) & (df['Time'] <= end_date)]
            
    #         if df_filtered.empty:
    #             raise ValueError("No hay datos en el rango de fechas proporcionado.")
            
    #         df_filtered = df_filtered[df_filtered['Message'] == 'Ended not OK']
    #         df_filtered = df_filtered[~df_filtered['Application'].isin(applications_to_exclude)]
    #         df_filtered['Notes'] = df_filtered['Notes'].str.replace(r'Alerta asociada al', 'Job errado', regex=True)
            
    #         return df_filtered
        
    #     except Exception as e:
    #         raise ValueError(f"Error procesando el archivo: {e}")



    def convert_reorder_and_filter(self, file_path, start_date, end_date):
        try:
            # Lista de aplicaciones a excluir
            applications_to_exclude = ['BDA', 'CALIDAD-RED', 'CONTROL-M', 'REDOYM', 'SEGURIDAD', 'SYM']

            # Orden esperado de las columnas en el DataFrame
            column_order = [
                'Run Counter', 'Origin', 'Severity', 'Status', 'Notes', 'Job Name', 'Time',
                'Message', 'Member/File Name', 'Application', 'Sub Application',
                'Control-M Server', 'Alert ID', 'Host ID', 'Order ID', 'Changed By',
                'Update Time', 'Run As', 'Remedy Ticket', 'Closed From EM'
            ]

            # Seleccionar el motor adecuado según la extensión del archivo
            if file_path.endswith('.xls'):
                df = pd.read_excel(file_path, engine='xlrd')
            elif file_path.endswith('.xlsx'):
                df = pd.read_excel(file_path, engine='openpyxl')
            else:
                raise ValueError("Formato de archivo no soportado")

            # Verificar que todas las columnas requeridas están presentes
            required_columns = set(column_order)
            if not required_columns.issubset(df.columns):
                missing_columns = required_columns - set(df.columns)
                raise ValueError(f"Columnas faltantes: {missing_columns}")

            # Reordenar las columnas según el orden esperado
            df = df[column_order]

            # Identificar columnas de texto y rellenar valores nulos con cadenas vacías
            text_columns = df.select_dtypes(include=['object', 'string']).columns
            df[text_columns] = df[text_columns].fillna('')

            # Convertir la columna 'Time' a formato datetime
            df['Time'] = pd.to_datetime(df['Time'], errors='coerce')

            # Verificar si la columna 'Time' contiene valores válidos
            if df['Time'].isna().all():
                raise ValueError("La columna 'Time' no contiene valores válidos.")

            # Convertir las fechas de inicio y fin a objetos datetime
            start_date = pd.to_datetime(start_date)
            end_date = pd.to_datetime(end_date + ' 23:59:59')

            # Filtrar por rango de fechas
            df_filtered = df[(df['Time'] >= start_date) & (df['Time'] <= end_date)]

            # Verificar si hay datos en el rango de fechas
            if df_filtered.empty:
                raise ValueError("No hay datos en el rango de fechas proporcionado.")

            # Filtrar filas donde 'Message' sea 'Ended not OK'
            df_filtered = df_filtered[df_filtered['Message'] == 'Ended not OK']

            # Excluir filas donde 'Application' esté en la lista de exclusión
            df_filtered = df_filtered[~df_filtered['Application'].isin(applications_to_exclude)]

            # Reemplazar texto en la columna 'Notes'
            df_filtered['Notes'] = df_filtered['Notes'].str.replace(r'Alerta asociada al', 'Job errado', regex=True)

            # Devolver el DataFrame filtrado
            return df_filtered

        except Exception as e:
            # Manejo de errores y descripción del problema
            raise ValueError(f"Error procesando el archivo: {e}")


    # Función para borrar todos los archivos y carpetas en una ruta dada excepto un archivo específico
    def limpiar_carpeta(self,carpeta, archivo_a_conservar):
        for root, dirs, files in os.walk(carpeta):
            for file in files:
                file_path = os.path.join(root, file)
                if file != archivo_a_conservar:
                    os.remove(file_path)
            for dir in dirs:
                dir_path = os.path.join(root, dir)
                os.rmdir(dir_path)

    def combinar_archivos_txt(self,archivos_txt, ruta_archivo_combinado):
        with open(ruta_archivo_combinado, 'w', encoding='utf-8') as archivo_combinado:
            for ruta_archivo in archivos_txt:
                with open(ruta_archivo, 'r', encoding='utf-8') as archivo:
                    contenido = archivo.read()
                    archivo_combinado.write(contenido)
                    archivo_combinado.write("\n" + "="*70 + "\n")

    def obtener_fecha_de_archivo(self,nombre_archivo):
        match = re.search(r'(\d{2})(\d{2})(\d{4})', nombre_archivo)
        if match:
            dia, mes, año = match.groups()
            return f"{dia}/{mes}/{año}"
        return ""

    def procesar_archivo(self,ruta_archivo, fecha):
        with open(ruta_archivo, 'r', encoding='utf-8') as file:
            contenido = file.read()
        
        secciones = re.split(r'={20,}', contenido)
        secciones_con_fecha = [f"{fecha}\n{seccion.strip()}" for seccion in secciones if seccion.strip()]
        nuevo_contenido = "\n" + "\n".join([f"{sec}\n{'='*70}" for sec in secciones_con_fecha])
        
        with open(ruta_archivo, 'w', encoding='utf-8') as file:
            file.write(nuevo_contenido)

    def extraer_txt_de_zip(self,ruta_zip, ruta_destino, ruta_archivo_combinado):
        if not os.path.exists(ruta_destino):
            os.makedirs(ruta_destino)
        
        archivos_txt = []
        
        # with zipfile.ZipFile(self,ruta_zip, 'r') as zip_ref:
        namelist = glob(os.path.join(ruta_destino,'*.txt'))

        for file in namelist:
            # print("file...................",file)
            # zip_ref.extract(file, ruta_destino)

            ruta_archivo = os.path.join(ruta_destino, file)

            archivos_txt.append(ruta_archivo)
            fecha = self.obtener_fecha_de_archivo(file)
            self.procesar_archivo(ruta_archivo, fecha)
        
        self.combinar_archivos_txt(archivos_txt, ruta_archivo_combinado)



    # def guardar_excel_con_formato(self, df, ruta_excel_salida):
    #     # Eliminar la columna 'INC' antes de guardar
    #     if 'INC' in df.columns:
    #         df.drop(columns=['INC'], inplace=True)
        
    #     # Crear un libro de trabajo de Excel
    #     wb = Workbook()
    #     ws = wb.active
        
    #     # Agregar el DataFrame a la hoja de Excel
    #     for r in dataframe_to_rows(df, index=False, header=True):
    #         ws.append(r)
        
    #     # Definir el color de relleno celeste para los encabezados
    #     celeste_fill = PatternFill(start_color="ADD8E6", end_color="ADD8E6", fill_type="solid")
        
    #     # Definir el color de relleno amarillo para la condición en las columnas 'Estado' y 'Notes Corregido'
    #     amarillo_fill = PatternFill(start_color="FFFF00", end_color="FFFF00", fill_type="solid")
        
    #     # Identificar las columnas de interés
    #     col_estado = None
    #     col_notes_corregido = None
        
    #     # Encontrar los índices de las columnas 'Estado' y 'Notes Corregido' en el encabezado
    #     for cell in ws[1]:  # Encabezado
    #         if cell.value == "Estado":
    #             col_estado = cell.column
    #             cell.fill = celeste_fill  # Pintar el título de la columna 'Estado'
    #         if cell.value == "Notes Corregido":
    #             col_notes_corregido = cell.column
    #             cell.fill = celeste_fill  # Pintar el título de la columna 'Notes Corregido'
        
    #     # Aplicar formato amarillo en las columnas 'Estado' y 'Notes Corregido' cuando el estado sea "CORREGIR EN CM"
    #     if col_estado and col_notes_corregido:
    #         for row in ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=col_estado, max_col=col_estado):
    #             for celda_estado in row:
    #                 if celda_estado.value and str(celda_estado.value).strip().upper() == "CORREGIR EN CM":
    #                     # Pintar la celda en la columna 'Estado' de amarillo
    #                     celda_estado.fill = amarillo_fill
    #                     # Pintar la celda correspondiente en 'Notes Corregido' de amarillo
    #                     corresponding_cell = ws.cell(row=celda_estado.row, column=col_notes_corregido)
    #                     corresponding_cell.fill = amarillo_fill
        
    #     # Agregar un filtro a todas las columnas del DataFrame
    #     ws.auto_filter.ref = ws.dimensions  # Aplica el filtro al rango total de datos (incluye encabezado)
        
    #     # Aplicar un filtro específico en la columna 'Estado' con el valor "CORREGIR EN CM"
    #     # Aunque `openpyxl` no permite fijar el filtro por defecto, Excel lo activará visualmente.
    #     if col_estado:
    #         estado_letter = get_column_letter(col_estado)
    #         ws.auto_filter.add_filter_column(col_estado - 1, ["CORREGIR EN CM"])  # Filtro por la palabra

    #     # Ajustar automáticamente el ancho de las columnas según el contenido
    #     for col in ws.columns:
    #         max_length = 0
    #         column = col[0].column_letter  # Obtener la letra de la columna
    #         for cell in col:
    #             try:
    #                 if len(str(cell.value)) > max_length:
    #                     max_length = len(cell.value)
    #             except:
    #                 pass
    #         adjusted_width = max_length + 2  # Agregar un pequeño margen
    #         ws.column_dimensions[column].width = adjusted_width
        
    #     # Guardar el archivo de Excel
    #     wb.save(ruta_excel_salida)



    def guardar_excel_con_formato(self, df, ruta_excel_salida):
        # Eliminar la columna 'INC' antes de guardar
        if 'INC' in df.columns:
            df.drop(columns=['INC'], inplace=True)
        
        # Ordenar el DataFrame por la columna 'Estado', colocando "CORREGIR EN CM" primero
        df['Estado'] = df['Estado'].fillna('')  # Rellenar valores nulos
        df['prioridad'] = df['Estado'].apply(lambda x: 0 if str(x).strip().upper() == 'CORREGIR EN CM' else 1)
        df = df.sort_values(by='prioridad', ascending=True).drop(columns=['prioridad'])  # Ordenar por la prioridad

        # Crear un libro de trabajo de Excel
        wb = Workbook()
        ws = wb.active
        
        # Agregar el DataFrame a la hoja de Excel
        for r in dataframe_to_rows(df, index=False, header=True):
            ws.append(r)
        
        # Definir el color de relleno celeste para los encabezados
        celeste_fill = PatternFill(start_color="ADD8E6", end_color="ADD8E6", fill_type="solid")
        
        # Definir el color de relleno amarillo para la condición en las columnas 'Estado' y 'Notes Corregido'
        amarillo_fill = PatternFill(start_color="FFFF00", end_color="FFFF00", fill_type="solid")
        
        # Identificar las columnas de interés
        col_estado = None
        col_notes_corregido = None
        
        # Encontrar los índices de las columnas 'Estado' y 'Notes Corregido' en el encabezado
        for cell in ws[1]:  # Encabezado
            if cell.value == "Estado":
                col_estado = cell.column
                cell.fill = celeste_fill  # Pintar el título de la columna 'Estado'
            if cell.value == "Notes Corregido":
                col_notes_corregido = cell.column
                cell.fill = celeste_fill  # Pintar el título de la columna 'Notes Corregido'
        
        # Aplicar formato amarillo en las columnas 'Estado' y 'Notes Corregido' cuando el estado sea "CORREGIR EN CM"
        if col_estado and col_notes_corregido:
            for row in ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=col_estado, max_col=col_estado):
                for celda_estado in row:
                    if celda_estado.value and str(celda_estado.value).strip().upper() == "CORREGIR EN CM":
                        # Pintar la celda en la columna 'Estado' de amarillo
                        celda_estado.fill = amarillo_fill
                        # Pintar la celda correspondiente en 'Notes Corregido' de amarillo
                        corresponding_cell = ws.cell(row=celda_estado.row, column=col_notes_corregido)
                        corresponding_cell.fill = amarillo_fill
        
        # Agregar un filtro a todas las columnas del DataFrame
        ws.auto_filter.ref = ws.dimensions  # Aplica el filtro al rango total de datos (incluye encabezado)
        
        # Ajustar automáticamente el ancho de las columnas según el contenido
        for col in ws.columns:
            max_length = 0
            column = col[0].column_letter  # Obtener la letra de la columna
            for cell in col:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(cell.value)
                except:
                    pass
            adjusted_width = max_length + 2  # Agregar un pequeño margen
            ws.column_dimensions[column].width = adjusted_width
        
        # Guardar el archivo de Excel
        wb.save(ruta_excel_salida)



    def comparar_inc_notes_only(self, row):
        # print("comparar_inc_notes_only....................",row)
        inc = row['INC']
        notes = row['Notes']
        # Asegurarse de que 'INC' y 'Notes' sean cadenas de texto
        if not isinstance(inc, str):
            inc = str(inc) if not pd.isna(inc) else ''
        if not isinstance(notes, str):
            notes = str(notes) if not pd.isna(notes) else ''

        print("inc...............",inc)
        print("notes...............",notes)

        match = re.search(r'INC\d{12}', notes)
        if match:
            inc_notes = match.group(0)
            print("inc_notes....................",inc_notes)
            if inc != inc_notes:
                nueva_notes = re.sub(r'INC\d{12}', inc, notes)
                return nueva_notes
            else:
                return notes
        else:

            # Validar si hay más de 12 dígitos después de 'INC'
            extra_digits_match = re.search(r'INC(\d{12,}|.*[A-Za-z].*)', notes)
            print("extra_digits_match...............",extra_digits_match)

            if extra_digits_match:
                inc_notes = extra_digits_match.group(0)
                nueva_notes = re.sub(r'INC(\d{12,}|.*[A-Za-z].*)', inc, notes)
                return nueva_notes + " "+"Tiene más de 12 dígitos (entre #, A-Z); corregir el txt."
                        
            if notes.strip() == "":
                print("Notes vacío, retornando como está.")
                return notes
            
            print("if notes...............", notes)
            return notes
        

    def comparar_inc_notes_estado(self, row):
        inc = row['INC']
        notes = row['Notes']

        # Asegurarse de que 'INC' y 'Notes' sean cadenas de texto
        if not isinstance(inc, str):
            inc = str(inc) if not pd.isna(inc) else ''
        if not isinstance(notes, str):
            notes = str(notes) if not pd.isna(notes) else ''
        
        match = re.search(r'INC\d{12}', notes)
        if match:
            inc_notes = match.group(0)
            print("inc_notes..........",inc_notes)
            return "OK" if inc == inc_notes else "CORREGIR EN CM"
        
        return "CORREGIR EN CM"

    def leer_excel_y_procesar_txt(self,ruta_excel, ruta_txt_combinado, ruta_excel_salida):
        df = pd.read_excel(ruta_excel)

        with open(ruta_txt_combinado, 'r', encoding='utf-8') as file:
            contenido_txt = file.read()
        
        secciones = re.split(r'\n={20,}\n', contenido_txt)


        df['INC'] = df.apply(lambda row: self.buscar_inc_en_txt(row, secciones), axis=1)        
        df['Estado'] = df.apply(self.comparar_inc_notes_estado, axis=1)
        # df['Notes'] = df.apply(self.comparar_inc_notes_only, axis=1)
        df['Notes Corregido'] = df.apply(self.comparar_inc_notes_only, axis=1)

        column_order = [
            'Run Counter', 'Origin', 'Severity', 'Status', 'Notes','Notes Corregido', 'INC', 'Estado', 'Job Name', 'Time',
            'Message', 'Member/File Name', 'Application', 'Sub Application', 'Control-M Server', 'Alert ID',
            'Host ID', 'Order ID', 'Changed By', 'Update Time', 'Run As', 'Remedy Ticket', 'Closed From EM'
        ]
        df = df.reindex(columns=column_order)

        # Condiciones para eliminar filas
        valores_application = ["CALIDAD-RED", "BDA", "REDOYM", "SYM", "SEGURIDAD", "CONTROLM", "BD"]
        valores_sub_application = ["OFS", "ETL_CC_PO", "ETL_CC_GI", "ETL_TOA", "GEODIR"]

        df = df[~df['Application'].isin(valores_application)]
        df = df[~((df['Application'] == 'SOAP') & (df['Sub Application'].isin(valores_sub_application)))]

        self.guardar_excel_con_formato(df, ruta_excel_salida)


    def buscar_en_seccion(self,termino, seccion):
        termino_lower = termino.lower()
        print("termino_lower...",termino_lower)
        seccion_lower = seccion.lower()
        print("seccion_lower...",seccion_lower)

        val=re.search(r'\b' + re.escape(termino_lower) + r'\b', seccion_lower) is not None

        print("val...............",val)
        return val
    

    def buscar_en_seccion_con_numeros(self,termino, seccion):
        # Convertir a minúsculas
        termino_lower = termino.lower()
        seccion_lower = seccion.lower()
        print("buscar_en_seccion_con_numeros ----- termino_lower..............",termino_lower)
        print("seccion_lower..............",seccion_lower)
        tiene_numeros_seccion =self.dividir_y_validar_numeros(seccion_lower)

        # si existe numeros en sección 
        if tiene_numeros_seccion:
            # Intentar buscar sin quitar los números
            val = re.search(r'\b' + re.escape(termino_lower) + r'\b', seccion_lower) is not None
            print("Resultado sin quitar números:", val)
        else:
            # Intentar buscar sin quitar los números
            val = re.search(r'\b' + re.escape(termino_lower) + r'\b', seccion_lower) is not None
            print("Resultado sin quitar números:", val)
            if not val:
                termino_lower_sin_numero = re.split(r'\d', termino_lower, 1)[0]
                print("Sección sin números........................................................................................:", termino_lower_sin_numero)
                val = re.search(termino_lower_sin_numero, seccion_lower) is not None
                print("Resultado quitando números:", val)
        return val
    
            

    def dividir_y_validar_numeros(self, texto):
        # Buscar la palabra "errado"
        match = re.search(r'\berrado\b', texto)
        if match:
            # Si se encuentra "errado", dividir el texto 6 caracteres antes
            texto_dividido = texto[:max(0, match.start() - 6)]
        else:
            # Si no se encuentra "errado", mantener el texto tal cual
            texto_dividido = texto
        texto_sin_fecha = re.sub(r'\d{2}/\d{2}/\d{4}\s+', '', texto_dividido)
        print("Texto dividido:---------------", texto_sin_fecha,'------------------fin ')

        # Validar si existen números en el texto dividido
        tiene_numeros = bool(re.search(r'\d', texto_sin_fecha))
        # print("¿El texto contiene números?", tiene_numeros)
        return tiene_numeros ,





    # def buscar_inc_en_txt(self, row, secciones):
    #     try:
    #         fecha = pd.to_datetime(row['Time']).strftime('%d/%m/%Y')
    #         application = row.get('Application', '')
    #         job_name = row.get('Job Name', '')
    #         sub_application = row.get('Sub Application', '')

    #         job_name_sin_numero = re.split(r'\d', job_name, 1)[0] if job_name else ''


    #             # (job_name, self.buscar_en_seccion_con_numeros),
    #             # (job_name_sin_numero, self.buscar_en_seccion_con_numeros),
    #         # Lista de términos a buscar
    #         terminos_busqueda = [
    #             (job_name, self.buscar_en_seccion),
    #             (job_name_sin_numero, self.buscar_en_seccion),
    #             (application, self.buscar_en_seccion),
    #             (sub_application, self.buscar_en_seccion)
    #         ]
            
    #         # Iterar sobre las secciones
    #         for seccion in secciones:
    #             if fecha not in seccion:
    #                 continue
    #             for termino, funcion_busqueda in terminos_busqueda:
    #                 if termino and funcion_busqueda(termino, seccion):
    #                     return self.extraer_inc(seccion)
            
    #     except KeyError as e:
    #         return f"Error: Campo faltante en la fila: {str(e)}"
    #     except Exception as e:
    #         return f"Error: {str(e)}"

    #     return "no se encuentra"




    def buscar_inc_en_txt(self,row, secciones):
        fecha = pd.to_datetime(row['Time']).strftime('%d/%m/%Y')
        
        application = row['Application']
        job_name = row.get('Job Name', '')
        sub_application = row.get('Sub Application', '')
        
        job_name_sin_numero = re.split(r'\d', job_name, 1)[0]






        print("primero job_name.............",job_name)
        # Verificar que hay un nombre de trabajo
        if job_name:
            for seccion in secciones:
                # Comprobar si LITCLU0085 está en la sección
                if job_name in seccion:
                    if fecha in seccion and self.buscar_en_seccion(job_name, seccion):
                    # if fecha in seccion and self.buscar_en_seccion_con_numeros(job_name, seccion):
                        return self.extraer_inc(seccion)
                        
        print("job_name_sin_numero. paso............",job_name_sin_numero)

        # Verificar que hay un nombre de trabajo
        if job_name_sin_numero:
            for seccion in secciones:
                # Comprobar si LITCLU está en la sección
                if job_name_sin_numero in seccion:
                    if fecha in seccion and self.buscar_en_seccion(job_name_sin_numero, seccion):
                        
                        return self.extraer_inc(seccion)

        if application:
            for seccion in secciones:
                if fecha in seccion and self.buscar_en_seccion(application, seccion):
                    return self.extraer_inc(seccion)
 
        if sub_application:
            for seccion in secciones:
                if fecha in seccion and self.buscar_en_seccion(sub_application, seccion):
                    return self.extraer_inc(seccion)
        
        return "no se encuentra"


    def extraer_inc(self, seccion):
        incidencia_a_ignorar = "INC000002574850"  # Define la incidencia que quieres ignorar
        
        match = re.search(r'INC\d{12}', seccion)
        if match:
            incidencia = match.group(0)
            print("incidencia..",incidencia)
            if incidencia == incidencia_a_ignorar:
                return "no se encuentra"
            return incidencia
        return "no se encuentra"


    def eliminar_carpeta_y_contenido(self,ruta_carpeta):
        if os.path.exists(ruta_carpeta):
            try:
                shutil.rmtree(ruta_carpeta)
                print(f'Carpeta eliminada: {ruta_carpeta}')
            except Exception as e:
                print(f'No se pudo eliminar {ruta_carpeta}. Razón: {e}')
        else:
            print(f'La ruta especificada no existe: {ruta_carpeta}')

        
class ReporteC7(viewsets.ModelViewSet):
    permission_classes = [IsAuthenticated]

    @action(detail=False, methods=['post'])
    def create(self, request, *args, **kwargs):
        datos_totales = request.data  
        archivo_excels = request.FILES.get('archivo', None)
        print("archivo_excels............................",archivo_excels)
        funcionesbasicasfiltroexcel = FuncionesBasicasExcelFiltros()
        ## Analisi excels
        carpeta_destino_base = r'C7'  
        nombre_mes_actual = datetime.now().strftime('%B').upper()
        codigo_aleatorio = funcionesbasicasfiltroexcel.generar_codigo_aleatorio()
        carpeta_destino = os.path.join(settings.MEDIA_ROOT,carpeta_destino_base, nombre_mes_actual+'_'+codigo_aleatorio)
        print("carpeta_destino..................",carpeta_destino)
        funcionesbasicasfiltroexcel.extraer_archivos(archivo_excels,carpeta_destino)
        # xls_files = glob(os.path.join(carpeta_destino,'*.xls'))

        # Buscar archivos con extensiones .xls y .xlsx
        xls_files = glob(os.path.join(carpeta_destino, '*.xls')) + glob(os.path.join(carpeta_destino, '*.xlsx'))
        # Opcional: ordenar los archivos
        xls_files.sort()

        print("xls_files----------------",xls_files)
        dataframes = []
        encontrado_archivo_invalido = False
        #rango de fechas
        fechas_ini=datos_totales.pop('fechas_ini', None)
        fechas_fin=datos_totales.pop('fechas_fin', None)


        start_date = fechas_ini[0]
        end_date = fechas_fin[0]

        print("start_date...........................",start_date)
        print("end_date..........................",end_date)

        for file_path in xls_files:
            df_filtered = funcionesbasicasfiltroexcel.convert_reorder_and_filter(file_path, start_date, end_date)

            if df_filtered is not None:
                dataframes.append(df_filtered)
            else:
                encontrado_archivo_invalido = True
                print(f"No se pudo procesar el archivo: {file_path}. Verifique las columnas requeridas.")
                if os.path.exists(carpeta_destino):
                    shutil.rmtree(carpeta_destino)
                    print("Carpeta eliminada:", carpeta_destino)
                else:
                    print("Directory does not exist:", carpeta_destino)
                return Response({"msg": "ERROR", "detail":f"No se pudo procesar el archivo:</br>  {file_path}. </br> Verifique las columnas requeridas." }, status=status.HTTP_400_BAD_REQUEST)
        
        if encontrado_archivo_invalido:
            print("Se encontraron archivos inválidos. Deteniendo el proceso.")
            if os.path.exists(carpeta_destino):
                shutil.rmtree(carpeta_destino)
                print("Carpeta eliminada:", carpeta_destino)
            else:
                print("Directory does not exist:", carpeta_destino)
            return Response({"msg": "ERROR", "detail":f"Se encontraron archivos inválidos.<br>  Deteniendo el proceso." }, status=status.HTTP_400_BAD_REQUEST)
        else:

            print("combined_df................---------------",dataframes)
            # Combinar los DataFrames en uno solo
            combined_df = pd.concat(dataframes, ignore_index=True)

            # Ordenar por 'Alert ID', luego por 'Notes' no nulos (descendente) y finalmente por 'Time' (descendente)
            combined_df.sort_values(by=['Alert ID', 'Notes', 'Time'], ascending=[True, False, False], inplace=True)

            # Eliminar duplicados de 'Alert ID', manteniendo la primera fila (con la prioridad dada por la ordenación)
            combined_df = combined_df.drop_duplicates(subset='Alert ID', keep='first')

            # Guardar el DataFrame combinado en un solo archivo Excel
            output_file = 'C7.xlsx'
            output_path = os.path.join(settings.MEDIA_ROOT, carpeta_destino, output_file)

            # Guardar el archivo Excel usando 'openpyxl' como motor
            combined_df.to_excel(output_path, index=False, engine='openpyxl')

            # etapa txts
            archivo_txt = request.FILES.get('archivotxt', None)
            codigo_aleatorio = funcionesbasicasfiltroexcel.generar_codigo_aleatorio()
            carpeta_destino_txt = os.path.join(settings.MEDIA_ROOT,carpeta_destino_base, nombre_mes_actual+'_'+codigo_aleatorio)# Subcarpeta JUNIO dentro de C7
            funcionesbasicasfiltroexcel.extraer_archivos(archivo_txt,carpeta_destino_txt)
            ruta_destino = os.path.join(carpeta_destino_txt)
            ruta_zip = ruta_destino+'/'+archivo_txt.name
            ruta_archivo_combinado = ruta_destino +'/combinado.txt'
            data_simulada=""
            with open(ruta_archivo_combinado, "w") as archivo_txt_combinado:
                archivo_txt_combinado.write(data_simulada)
            ruta_excel = os.path.join(settings.MEDIA_ROOT,carpeta_destino, output_file)
            ruta_excel_salida = os.path.join(settings.MEDIA_ROOT,carpeta_destino, output_file)
            funcionesbasicasfiltroexcel.extraer_txt_de_zip(ruta_zip, ruta_destino, ruta_archivo_combinado)

            
            funcionesbasicasfiltroexcel.leer_excel_y_procesar_txt(ruta_excel, ruta_archivo_combinado, ruta_excel_salida)
            print(f"Se ha limpiado la carpeta, dejando solo el archivo '{output_file}'")


            ruta_archivo = os.path.join(settings.MEDIA_ROOT, carpeta_destino, output_file)
            ruta_app = "/usr/src/app/"
            ruta_local = "http://localhost:8000/"
            nueva_ruta = ruta_archivo.replace(ruta_app, ruta_local)
            
            # Remove the directory if it exists
            if os.path.exists(carpeta_destino_txt):
                shutil.rmtree(carpeta_destino_txt)
                print("carpeta txt:", carpeta_destino_txt)
            else:
                print("carpeta txt does not exist:", carpeta_destino_txt)
        return Response(nueva_ruta, status=status.HTTP_200_OK)
        
    @action(detail=False, methods=['put'])
    def update(self, request, *args, **kwargs):
        partial = kwargs.pop('partial', False)
        datos_totales = request.data  # Make a mutable copy of the request data

        ruta_directorio = datos_totales.get('ruta_directorio')

        if ruta_directorio:
            print("ruta_directorio:", ruta_directorio)
        else:
            print("ruta_directorio not found in the request data.")
            return Response({"error": "ruta_directorio not found"}, status=400)

        dir_path = ruta_directorio.replace('http://localhost:8000/media/', '')
        full_dir_path = os.path.join(settings.MEDIA_ROOT, dir_path)

        # Print the directory path to be removed
        print("Directory to be removed:", full_dir_path)

        # Remove the directory if it exists
        if os.path.exists(full_dir_path):
            shutil.rmtree(full_dir_path)
            print("Directory removed:", full_dir_path)
            return Response({"message": "Directory removed successfully"})
        else:
            print("Directory does not exist:", full_dir_path)
            return Response({"error": "Directory does not exist"}, status=404)
        
#boot
        
class ChatbootViewSet(viewsets.ModelViewSet):
    permission_classes = [IsAuthenticated]

    def enviar_mensaje_usu(self, request):
        user_message = request.data.get('mensaje_usu')
        bot_last_message = request.data.get('mensaje_bot')
        
        
        print("user_message.......................",user_message)
        print("bot_last_message.......................",bot_last_message)
        if not user_message:
            return Response({'error': 'No se envió ningún mensaje'}, status=status.HTTP_400_BAD_REQUEST)

        # bot_response = self.chatbot.get_response(user_message)
        # return Response({'response': bot_response.text}, status=status.HTTP_200_OK)
        # return Response({'response': 'respuesta del api'}, status=status.HTTP_200_OK)
        bot_response = get_bot_response(user_message, bot_last_message)

        # if 'error' in bot_response:
        #     return Response({'error': 'No se envió ningún mensaje'}, status=status.HTTP_400_BAD_REQUEST)
        # else:
        return JsonResponse({'response': bot_response})

    # def send_message(request):
    #     user_message = request.GET.get('message', '')
    #     bot_response = get_bot_response(user_message)
    #     return JsonResponse({'response': str(bot_response)})



class LineaTiempoViewSet(viewsets.ModelViewSet):
    queryset = LineaTiempo.objects.all()
    serializer_class = LineaTiempoSerializer
    permission_classes = [IsAuthenticated]

    @action(detail=False, methods=['get'])
    def list(self, request, *args, **kwargs):
        queryset = self.get_queryset()
        querysetado = self.queryset.filter(status=True)
        serializer = self.get_serializer(querysetado, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)
    
    @action(detail=False, methods=['get'])
    def list_id(self, request,pk):
        datas = self.queryset.filter(id=pk , status=True)
        serializer = self.get_serializer(datas, many=True)
        return Response(serializer.data)    
    
    @action(detail=False, methods=['get'])
    def list_by_type(self, request, id, type=None):
        if type == 'usuario':
            archivos = self.queryset.filter(usuarios=id,status=True)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)

    @action(detail=False, methods=['post'])
    def create(self, request, *args, **kwargs):
        datos_totales = request.data
        print("datos_totales-........................",datos_totales)
        return Response(datos_totales, status=status.HTTP_200_OK)
        
    @action(detail=False, methods=['put'])
    def update(self, request, *args, **kwargs):
        partial = kwargs.pop('partial', False)
        datos_totales = request.data
        return Response(datos_totales, status=status.HTTP_200_OK)

class LineaTiempoEventoViewSet(viewsets.ModelViewSet):

    queryset = LineaTiempoEvento.objects.all()
    serializer_class = LineaTiempoEventoSerializer
    permission_classes = [IsAuthenticated]


    @action(detail=False, methods=['get'])
    def list(self, request, *args, **kwargs):
        queryset = self.get_queryset()
        print("queryset lin_tiemp",queryset)
        querysetado = self.queryset.filter(status=True)
        serializer = self.get_serializer(querysetado, many=True)

        return Response(serializer.data, status=status.HTTP_200_OK)
    
    @action(detail=False, methods=['get'])
    def list_id(self, request,pk):

        datas = self.queryset.filter(id=pk , status=True)
        serializer = self.get_serializer(datas, many=True)
        return Response(serializer.data)    
    
    @action(detail=False, methods=['get'])
    def list_by_type(self, request, id, type=None):
        if type == 'usuario':
            archivos = self.queryset.filter(usuarios=id,status=True)
        else:
            return Response({"detail": "Invalid type"}, status=400)
        
        serializer = self.get_serializer(archivos, many=True)
        return Response(serializer.data)

    @action(detail=False, methods=['post'])
    def create(self, request, *args, **kwargs):
        datos_totales = request.data

        return Response(datos_totales, status=status.HTTP_200_OK)

        
    @action(detail=False, methods=['put'])
    def update(self, request, *args, **kwargs):
        partial = kwargs.pop('partial', False)
        datos_totales = request.data
        return Response(datos_totales, status=status.HTTP_200_OK)
